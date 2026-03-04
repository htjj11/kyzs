'''
此处定义所有的接口的方法
'''

import requests
from api.sql_role import kyzs_sql
from config import *
import json
import datetime
import os
import pypandoc
import base64
import hashlib
import time
import urllib.parse



#三种从oilink获取信息的方法
def get_articles_from_oillink(keywords_list:list,page:int,size:int,user_id:int):
    """
    oillink的文献检索接口
    执行 curl 请求，从 http://data.oillink.com 获取文章搜索结果
    
    返回:
        requests.Response 对象，包含请求的响应
        格式为：{'code': 200, 'msg': 'success', 'data': [{'abstract_zh': '直线型振动筛和...', 'id': '6860e2f5163c01c850d9987f', 'keywords': None, 'keywords_zh': ['钻井液振动筛', '双振型', '离散单元法', 'EDEM', '筛分效率'], 'title_zh': '基于离散单元法的双振型下钻井液振动筛不 同颗粒筛分效率对比研究 后印本', 'year': 2025}, {'abstract': 'The primary ..', 'abstract_zh': '的目的...', 'doi': '10.3969/j.issn.1000-6532.2025.02.012', 'id': '68108681163c01c850d8d66d', 'keywords': ['Mining processing engineering', 'Flotation', 'Barite', 'Waste drilling fluid', 'Inhibitors', 'Collectors', 'Mechanism study'], 'keywords_zh': ['矿物加工工程', '浮选', '重晶石', '废弃钻井液', '抑制剂', '捕收剂', '机理研究'], 'title': 'Flotation Recovery of Barite in Waste Drilling Fluids Using Novel Inhibitors', 'title_zh': '新型抑制剂对废弃钻井液中重晶石的浮选回收', 'year': 2025}]}
    """
    url = "http://data.oillink.com/api/shengli/articlesearch/index"
    params = {
        "keywords": f'{keywords_list}',
        "page": page,
        "size": size
    }
    try:
        response = requests.get(url, params=params)
        response.raise_for_status()  # 检查请求是否成功

        #处理返回数据总的data列表，添加一个字段is_collected,遍历其中的doi是否在收藏列表中，如果在则为1，否则为0
        data = response.json()['data']

        #如果返回的data为None,则直接返回空列表
        if data == None:
            return []

        for article in data:
            doi = article.get('doi')
            if not doi:
                continue
            sql_sentence = f"""
            SELECT COUNT(*) FROM knowledgebase WHERE mark_info LIKE '{doi}' AND type_id = 1 AND user_id = {user_id};
            """
            result = kyzs_sql.mysql_exec(sql_sentence)
            if result[0]['COUNT(*)'] > 0:
                article['is_collected'] = 1
            else:
                article['is_collected'] = 0    
        return data


    except requests.RequestException as e:
        print(f"请求出错: {e}")
        return None

def get_patents_from_oillink(query:str, page:int = 0, size:int = 5):
    """
    专利检索执行 POST 请求，从 http://data.oillink.com 获取专利搜索结果
    参数:
        query (str): 检索关键词
        page (int): 页码，默认为 0
        size (int): 每页数量，默认为 2
    
    返回:
        dict: 请求的响应数据，请求失败返回 None
        返回格式：
        {
        "data":[[            {
                "abstract": {
                    "en": [
                        "NOVELTY e ADV"
                    ],
                    "zh": [
                        "一种自动驾。"
                    ]
                },
                "app_date": {
                    "seconds": 1669593600
                },
                "app_num": "202280026881",
                "applicant": [
                    {
                        "name": "北京百度网讯科技有限公司",
                        "raw_address_info": "100085 北京市海淀区上地十街10号百度大厦2层",
                        "sequence": 1
                    },
                    {
                        "name": "阿波罗智行美国有限公司",
                        "sequence": 2
                    }
                ],
                "assignee": [
                    {
                        "name": "BEIJING BAIDU NETCOM SCI & TECHNOLOGY CO (BDCA-C)",
                        "sequence": 1
                    },
                    {
                        "name": "APOLLO ZHIXING USA CO LTD (APLO-C)",
                        "sequence": 2
                    },
                    {
                        "name": "APOLLO AUTONOMOUS DRIVING USA LLC (APLO-C)",
                        "sequence": 3
                    }
                ],
                "claims": {
                    "zh": [
                        "1.一种运。",
                        "2.根据权学生边缘模型。"
                    ]
                },
                "country": "cn",
                "description": {
                    "zh": [
                        "技术领域",
                        " 。"
                    ]
                },
                "id": "661f54c1503a87c91e32a626",
                "inventor": [
                    {
                        "name": "寇浩锋",
                        "sequence": 1
                    },
                    {
                        "name": "朱晓毅",
                        "sequence": 2
                    },
                    {
                        "name": "张满江",
                        "sequence": 3
                    },
                    {
                        "name": "潘海伦",
                        "sequence": 4
                    }
                ],
                "ipc": [
                    {
                        "l1": "G",
                        "l2": "G05",
                        "l3": "G05B",
                        "l4": "G05B019/042"
                    },
                    {
                        "l1": "G",
                        "l2": "G06",
                        "l3": "G06F",
                        "l4": "G06F009/451"
                    },
                    {
                        "l1": "G",
                        "l2": "G06",
                        "l3": "G06N",
                        "l4": "G06N020/00"
                    }
                ],
                "ipcr": [
                    {
                        "l1": "G",
                        "l2": "G05",
                        "l3": "G05B",
                        "l4": "G05B19/042"
                    }
                ],
                "pct": {
                    "app_num": "PCT/CN2022/134810",
                    "pub_num": null
                },
                "priority": [
                    {
                        "country": "cn",
                        "date": {
                            "seconds": 1669564800
                        },
                        "num": "80026881"
                    },
                    {
                        "country": "unset",
                        "date": {
                            "seconds": 1669564800
                        },
                        "num": "134810"
                    },
                    {
                        "country": "cn",
                        "date": {
                            "seconds": 1696608000
                        },
                        "num": "80026881"
                    }
                ],
                "pub_date": {
                    "seconds": 1708646400
                },
                "pub_kind": "A",
                "pub_num": "117597637",
                "title": {
                    "en": [
                        "Method for ."
                    ],
                    "zh": [
                        "人工智能车辆操作系统"
                    ]
                }
            }
        ]]]
        }
    """
    url = "http://data.oillink.com/api/shengli/patentsearch/index"
    params = {
        "query": query,
        "page": page,
        "size": size
    }
    try:
        response = requests.post(url, params=params)
        response.raise_for_status()  # 检查请求是否成功
        data = response.json()
        #判断code是否为200
        if data['code'] != 200:
            return []
        data= data['data']

        new_data = []
        for item in data:
            # item 是列表，取出字典
            i = item[0]    

            # 获取专利id
            patent_id = i['id']
            sql = f"SELECT COUNT(*) FROM knowledgebase WHERE mark_info='{patent_id}' AND type_id = 2"
            id_result = kyzs_sql.mysql_exec(sql)[0]['COUNT(*)']

            # 设置是否已存在标记
            i['is_collected'] = 1 if id_result else 0

            # 转换时间字段
            i['app_date'] = datetime.datetime.fromtimestamp(i['app_date']['seconds'])
            i['pub_date'] = datetime.datetime.fromtimestamp(i['pub_date']['seconds'])

            # 添加到新列表
            new_data.append(i)
        return new_data
    except requests.RequestException as e:
        print(f"请求出错: {e}")
        return None

def get_online_infomation_api(keyword:str):
    """
    获取讯飞互联网信息信息
    :param keyword: 关键词
    :return: 在线信息
    """
    #讯飞网络知识检索
    def xunfei_online_search(question):
        url = "https://spark-api-open.xf-yun.com/v1/chat/completions"
        print(f'方法内：讯飞网络知识检索关键词：{question}', flush=True)
        # 定义要使用的工具
        tools = [
            {
                "type": "function",
                "function": {
                    "name": "get_research_information",
                    "description": "获取互联网科研信息",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "title": {
                                "type": "string",
                                "description": "信息标题"
                            },
                            "content": {
                                "type": "string",
                                "description": "信息内容，约1000字，倾向于真实案例和数据"
                            },
                            "source": {
                                "type": "string",
                                "description": "信息来源"
                            },
                            "date": {
                                "type": "string",
                                "description": "信息日期，格式为YYYY-MM-DD"
                            },
                            "author": {
                                "type": "string",
                                "description": "信息作者"
                            }
                        },
                        "required": ["title", "content", "source", "date", "author"]
                    }
                }
            },
            {
                "type": "web_search",
                "web_search": {
                    "enable": True,
                    "show_ref_label": True,
                    "search_mode": "deep"  # deep:深度搜索 / normal:标准搜索
                }
            }
        ]

        data = {
            "model": "4.0Ultra",
            "user": "default_user",  # 默认用户ID
            "messages": [
                {
                "role": "system",
                "content": "你是知识渊博的助理，能够获取互联网科研信息"
            },
                {
                    "role": "user",
                    "content": question
                }
            ],
            "temperature": 1,
            "top_k": 6,
            "stream": False,
            "max_tokens": 5000,
            "tools": tools,
            "tool_choice": {
                "type": "function",
                "function": {
                    "name": "get_research_information"
                }
            },

        }
        header = {"Authorization": "Bearer fXMOutpFGqmwUoTTSpHz:OTsqdwPkqNmvlKJouKSU"}
        print(f'准备执行请求：{url}，请求头：{header}，请求数据：{data}', flush=True)
        try:
            response = requests.post(url, headers=header, json=data)
            print(f'请求已发送，等待响应...', flush=True)
        except Exception as e:
            print(f'发送请求时出错：{e}', flush=True)
            return None

        return response.json()


    return xunfei_online_search(keyword)

# def get_online_infomation_summary_api(online_infomation:str):
#     """
#     获取长城网络检索一篇摘要
#     :param online_infomation: 互联网信息
#     :return: 互联网信息摘要
#     """
#     def changcheng_online_summary_api(question):
#         """
#         使用长程在线摘要API
#         :param question: 需要摘要的内容
#         :return: 摘要结果
#         """
#         import requests
#         import json
        
#         # API配置
#         url = "https://agent.ai.sinopec.com/aicoapi/gateway/v2/workflow/api_run/30c0e0b0728c4d3ca8f3d6703e640ddc"
#         api_key = "ckMWXPiMCdCV3Xyopy8shCcSHXloHx22"  # 需要替换为实际的API密钥
        
#         # 请求头设置
#         headers = {
#             "Content-Type": "application/json",
#             "Authorization": f"Bearer {api_key}"
#         }
        
#         # 请求体设置
#         payload = {
#             "content": [
#                 {
#                     "field_name": "content",
#                     "type": "input",
#                     "value": question
#                 }
#             ],
#             "stream": False
#         }
        
#         try:
#             # 发送POST请求，设置stream=True以处理流式响应
#             response = requests.post(url, headers=headers, data=json.dumps(payload), stream=False,verify=False)
#             response.raise_for_status()  # 检查请求是否成功
#             return response.text
#         except requests.exceptions.RequestException as e:
#             print(f"请求错误: {str(e)}")
#             return f"摘要API请求失败: {str(e)}"
#         except Exception as e:
#             print(f"处理响应时出错: {str(e)}")
#             return f"处理摘要结果时出错: {str(e)}"

#     # 调用长程在线摘要API
#     answer = changcheng_online_summary_api(online_infomation)
#     #获取文本呈现

#     answer = json.loads(answer)['data']['data']['文本呈现']
#     answer = answer.replace('\n', '\\n')
#     print(answer)
    
#     answer = json.loads(answer) 
#     # answer = re.sub(r'[\x00-\x1f\x7f]', '', answer)
#     # import ast  # 用来安全地解析 Python 字面量
#     # answer = ast.literal_eval(answer)
#     # # 第二步：对列表中每个字符串再进行一次 JSON 解析
#     # answer = [json.loads(item) for item in answer]
#     # answer = json.dumps(answer, ensure_ascii=False, indent=2)
#     # 调用大模型
#     return answer


#重庆聚合文献接口
def get_article_from_juhe_api(keywords:list,date:str=None,page=1,size=10,sort=3):
    """
    从重庆聚合文献接口获取文献
    :param keywords: 关键词
    :param date: 查询年份
    :param page: 页码
    :param pagesize: 每页数量
    :return: 文献列表
    """

    def md5_encrypt(text: str) -> str:
        """MD5加密函数"""
        md5 = hashlib.md5()
        md5.update(text.encode('utf-8'))
        return md5.hexdigest()
    exps = []
    for keyword in keywords:
        exps.append(f"关键词:{keyword}")
    if date:
        exps.append(f"年份:{date}")
    exp = " AND ".join(exps)
    exp = urllib.parse.quote(exp)
    params = "{}|{}|{}|{}".format(exp, page, size, sort)
    print('表达式：',params)
    times = "{}".format(int(time.time() * 1000))
    en_params = base64.b64encode(params.encode('utf-8')).decode('utf-8')
    sign = "{}|{}|{}".format(en_params, "AF85101E523744ADA233DF14CCC76980", times)
    en_sign = md5_encrypt(sign)
    # 设置请求头
    headers = {
        "User-Agent": "test",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection": "keep-alive",
    }
    # 构造请求参数
    query_params = {
        "params": en_params,
        "sign": en_sign,
        "token": "AF85101E523744ADA233DF14CCC76980",
        "times": times
    }
    try:
        # 发送带会话管理的请求
        with requests.Session() as session:
            # 首次请求获取可能的Cookie
            url = "http://61.128.134.70:6655/groups/ask/search"
            session.get(url, headers=headers, params=query_params, timeout=10)

            # 再次请求确保Cookie生效
            response = session.get(url, headers=headers, params=query_params, timeout=10)

            # 输出原始响应内容（用于调试）
            print("Raw Response:")
            print(response.text)

            # 检查HTTP状态码
            response.raise_for_status()
            # print("reach here!")
            # 解析JSON响应
            result = response.json()['result']
            # 处理业务逻辑
            paperlist = []
            for record in result['records']:
                paperlist.append(
                    {
                        "标题": f"{record['title']}",
                        "关键词": f"{record['keyword']}",
                        "年份": f"{record['year']}",
                        "摘要": f"{record['content']}",
                        "DOI": f"{record['doi']}",
                        "下载链接": f"http://61.128.134.70:6655{record['full_source']}"
                    }
                )
            return paperlist
    except requests.exceptions.RequestException as e:
        print(f"请求发生错误: {str(e)}")
    except ValueError as e:
        print(f"JSON解析失败，请检查响应内容是否为有效JSON: {str(e)}")
    except KeyError as e:
        print(f"响应数据缺少必要字段: {str(e)}")
    except Exception as e:
        print(f"未知错误: {str(e)}")

#钻井院万方接口
import requests
import xml.etree.ElementTree as ET
def get_article_from_wanfang_api(Datewithin = None, Keywords = [], StartRecord = 1, MaximumRecords = 10):
    exps = []

    for keyword in Keywords:
        exps.append(f"{keyword}")

    if Datewithin:
        exps.append(f'Date within "{Datewithin}-01-01 {Datewithin}-12-31"')
    exp = " and ".join(exps)

    print(f"检索表达式：{exp}")
    def send_soap_request_paper(exp, startRecord = 1, maximumRecords = 10):
        url = "http://10.68.16.2/S/SRW/Paper.asmx"
        headers = {
            "Host": "10.68.16.2",
            "Content-Type": "text/xml; charset=utf-8",
            "SOAPAction": "searchRetrieve"
        }
        
        soap_body = f"""<?xml version="1.0" encoding="utf-8"?>
                    <soap:Envelope xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance" 
                                xmlns:xsd="http://www.w3.org/2001/XMLSchema" 
                                xmlns:soap="http://schemas.xmlsoap.org/soap/envelope/">
                    <soap:Body>
                        <searchRetrieveRequest xmlns="http://www.loc.gov/zing/srw/">
                            <version>1.2</version>
                            <query>{exp}</query>
                            <operation>searchretrieve</operation>
                            <startRecord>{startRecord}</startRecord>
                            <maximumRecords>{maximumRecords}</maximumRecords>
                        </searchRetrieveRequest>
                    </soap:Body>
                    </soap:Envelope>"""
        try:
            response = requests.post(url, data=soap_body, headers=headers)
            response.raise_for_status()
            print("Response Status Code:", response.status_code)
            # print("Response Content:")
            # print(response.text)
            return response.text
        except requests.exceptions.RequestException as e:
            print("Error sending SOAP request:", e)
            return None
    response = send_soap_request_paper(exp, startRecord=StartRecord, maximumRecords=MaximumRecords)
    with open("raw_response.txt", "w") as f:
        f.write(f"{response}")
    # 解析响应
    def parse_academic_papers(soap_response):
        """
        解析SOAP响应，提取学术论文信息并生成结构化列表
        
        参数:
            soap_response: SOAP响应的XML字符串（str类型）
        返回:
            list: 每篇论文为一个字典，包含核心学术字段
        """
        # 1. 定义所有必要的命名空间映射（关键！没有则无法定位节点）
        namespaces = {
            "soap": "http://schemas.xmlsoap.org/soap/envelope/",  # SOAP信封命名空间
            "srw": "http://www.loc.gov/zing/srw/",               # 搜索响应命名空间
            "dc": "info:srw/schema/1/dc-v1.1",                  # 论文元数据（DC标准）命名空间
            "srw_dc": "info:srw/schema/1/dc-v1.1"               # srw_dc前缀命名空间
        }
        print(soap_response)
        # 2. 解析XML字符串，获取根节点
        root = ET.fromstring(soap_response)
        number_element = root.find('.//srw:numberOfRecords', namespaces=namespaces)

        # 提取数值
        if number_element is not None:
            number_of_records = number_element.text
            print(f"共解析到论文数量: {number_of_records}")  # 输出：numberOfRecords: 2506
        else:
            print("未找到numberOfRecords元素")
        # 3. 定位所有<record>节点（XPath路径：逐层穿透到records下的record）
        # 注：XPath中需用"命名空间:节点名"格式，且传入namespaces参数
        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record",
            namespaces=namespaces
        )

        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record",
            namespaces=namespaces
        )

        # 4. 遍历每个record，提取论文信息
        papers = []
        for idx, record in enumerate(record_nodes, 1):
            # 定位当前record下的论文元数据节点（recordData → srw_dc:dc）
            dc_node = record.find("srw:recordData/srw_dc:dc", namespaces=namespaces)
            # 提取单篇论文的核心字段（用.find()定位dc前缀的节点，text获取值，None时用空字符串兜底）
            paper_info = {
                # "序号": idx,  # 手动添加序号，方便查看
                "标题": dc_node.find("dc:title", namespaces).text.strip() if dc_node.find("dc:title", namespaces).text is not None else "",
                # "作者": dc_node.find("dc:Creator", namespaces).text.strip() if dc_node.find("dc:Creator", namespaces) is not None else "",
                "关键词": dc_node.find("dc:Subject", namespaces).text.strip() if dc_node.find("dc:Subject", namespaces).text is not None else "",
                "摘要": dc_node.find("dc:Description", namespaces).text.strip() if dc_node.find("dc:Description", namespaces).text is not None else "",
                # "贡献者": dc_node.find("dc:Contributor", namespaces).text.strip() if dc_node.find("dc:Contributor", namespaces) is not None else "",
                "发表时间": dc_node.find("dc:Date", namespaces).text.strip() if dc_node.find("dc:Date", namespaces).text is not None else "",
                # "文献类型": dc_node.find("dc:Type", namespaces).text.strip() if dc_node.find("dc:Type", namespaces) is not None else "",
                # "文件格式": dc_node.find("dc:Format", namespaces).text.strip() if dc_node.find("dc:Format", namespaces) is not None else "",
                "DOI": dc_node.find("dc:Identifier", namespaces).text.strip() if dc_node.find("dc:Identifier", namespaces).text is not None else "",
                # "语言": dc_node.find("dc:Language", namespaces).text.strip() if dc_node.find("dc:Language", namespaces) is not None else "",
                # "来源期刊": dc_node.find("dc:source", namespaces).text.strip() if dc_node.find("dc:source", namespaces) is not None else "",  # source首字母小写，与其他dc节点区分
                # "记录位置": record.find("srw:recordPosition", namespaces).text.strip() if record.find("srw:recordPosition", namespaces) is not None else ""
            }
            paper_info["发表时间"] = paper_info["发表时间"].split("-")[0]
            papers.append(paper_info)

        return (papers, number_of_records)

    return parse_academic_papers(response)
#钻井院万方专利接口
import requests
import xml.etree.ElementTree as ET
def wangfang_patent(Datewithin = None, patent_name = [], StartRecord = 1, MaximumRecords = 10):
    exps = []

    for keyword in patent_name:
        exps.append(f"{keyword}")
        
    if Datewithin:
        exps.append(f'F_PublicationDate within "{Datewithin}-01-01 {Datewithin}-12-31"')
    exp = " and ".join(exps)

    # print(f"检索表达式：{exp}")
    def send_soap_request_patent(exp, startRecord = 1, maximumRecords = 10):
        url = "http://10.68.16.2/S/SRW/Patent.asmx"
        headers = {
            "Host": "10.68.16.2",
            "Content-Type": "text/xml; charset=utf-8",
            "SOAPAction": "searchRetrieve"
        }
        
        soap_body = f"""<?xml version="1.0" encoding="utf-8"?>
                    <soap:Envelope xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance" 
                                xmlns:xsd="http://www.w3.org/2001/XMLSchema" 
                                xmlns:soap="http://schemas.xmlsoap.org/soap/envelope/">
                    <soap:Body>
                        <searchRetrieveRequest xmlns="http://www.loc.gov/zing/srw/">
                            <version>1.2</version>
                            <query>{exp}</query>
                            <operation>searchretrieve</operation>
                            <startRecord>{startRecord}</startRecord>
                            <maximumRecords>{maximumRecords}</maximumRecords>
                        </searchRetrieveRequest>
                    </soap:Body>
                    </soap:Envelope>"""
        try:
            response = requests.post(url, data=soap_body, headers=headers)
            response.raise_for_status()
            print("Response Status Code:", response.status_code)
            return response.text
        except requests.exceptions.RequestException as e:
            print("Error sending SOAP request:", e)
            return None

    response = send_soap_request_patent(exp, startRecord=StartRecord, maximumRecords=MaximumRecords)

    # 解析响应
    def parse_patents(soap_response) -> tuple:
        """
        解析SOAP响应，提取学术论文信息并生成结构化列表
        
        参数:
            soap_response: SOAP响应的XML字符串（str类型）
        返回:
            list: 每篇论文为一个字典，包含核心学术字段
        """
        # 1. 定义所有必要的命名空间映射（关键！没有则无法定位节点）
        namespaces = {
            "soap": "http://schemas.xmlsoap.org/soap/envelope/",  # SOAP信封命名空间
            "srw": "http://www.loc.gov/zing/srw/",               # 搜索响应命名空间
            "dc": "info:srw/schema/1/dc-v1.1",                  # 论文元数据（DC标准）命名空间
            "srw_dc": "info:srw/schema/1/dc-v1.1"               # srw_dc前缀命名空间
        }
        # print(soap_response)
        # 2. 解析XML字符串，获取根节点
        root = ET.fromstring(soap_response)
        number_element = root.find('.//srw:numberOfRecords', namespaces=namespaces)

        # 提取数值
        if number_element is not None:
            number_of_records = number_element.text
            print(f"共解析到论文数量: {number_of_records}")  # 输出：numberOfRecords: 2506
        else:
            print("未找到numberOfRecords元素")
        # 3. 定位所有<record>节点（XPath路径：逐层穿透到records下的record）
        # 注：XPath中需用"命名空间:节点名"格式，且传入namespaces参数
        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record",
            namespaces=namespaces
        )

        # 4. 遍历每个record，提取论文信息
        papers = []
        for idx, record in enumerate(record_nodes, 1):
            # 定位当前record下的论文元数据节点（recordData → srw_dc:dc）
            dc_node = record.find("srw:recordData/srw_dc:dc", namespaces=namespaces)

            # 提取单篇论文的核心字段（用.find()定位dc前缀的节点，text获取值，None时用空字符串兜底）
            paper_info = {
                "申请号": dc_node.find("dc:ApplicationNo", namespaces).text.strip() if dc_node.find("dc:ApplicationNo", namespaces).text is not None else "",
                "申请日": dc_node.find("dc:ApplicationDate", namespaces).text.strip() if dc_node.find("dc:ApplicationDate", namespaces).text is not None else "",
                "公开号": dc_node.find("dc:PublicationNo", namespaces).text.strip() if dc_node.find("dc:PublicationNo", namespaces).text is not None else "",
                "公开日": dc_node.find("dc:PublicationDate", namespaces).text.strip() if dc_node.find("dc:PublicationDate", namespaces).text is not None else "",
                "专利名称": dc_node.find("dc:PatentName", namespaces).text.strip() if dc_node.find("dc:PatentName", namespaces).text is not None else "",
                "申请人": dc_node.find("dc:Applicant", namespaces).text.strip() if dc_node.find("dc:Applicant", namespaces).text is not None else "",
                "发明人": dc_node.find("dc:Inventor", namespaces).text.strip() if dc_node.find("dc:Inventor", namespaces).text is not None else "",
                "IPC分类号": dc_node.find("dc:ClassMain", namespaces).text.strip() if dc_node.find("dc:ClassMain", namespaces).text is not None else "",
                "摘要": dc_node.find("dc:Abstract", namespaces).text.strip() if dc_node.find("dc:Abstract", namespaces).text is not None else "",
                "权利要求": dc_node.find("dc:SignoryItem", namespaces).text.strip() if dc_node.find("dc:SignoryItem", namespaces).text is not None else "",
            }
            papers.append(paper_info)

        return (papers, number_of_records)

    return parse_patents(response)
    

    
#大模型方法
def siliconflow_deepseek_answer(question):
    """
    deepseek大模型问答接口
    :param question: 题目
    :return:  回复文本（字符串）
    """
    url = "https://api.siliconflow.cn/v1/chat/completions"
    payload = {
        "model": "Pro/deepseek-ai/DeepSeek-V3",
        "messages": [
            {
                "role": "user",
                "content": str(question)
            }
        ],
        "stream": False,
        "max_tokens": 8192,
        "stop": ["null"],
        "temperature": 0.7,
        "top_p": 0.7,
        "top_k": 50,
        "frequency_penalty": 0.5,
        "n": 1,
        "response_format": {"type": "text"},
        "tools": [
            {
                "type": "function",
                "function": {
                    "description": "<string>",
                    "name": "<string>",
                    "parameters": {},
                    "strict": False
                }
            }
        ]
    }
    headers = {
        "Authorization": "Bearer sk-wmsgbfgsvjxjmyopswmaqfxnwtgmvtwqgsigehxmgwoihgeg",
        "Content-Type": "application/json"
    }

    response = requests.request("POST", url, json=payload, headers=headers)

    response = response.json()

    """
    返回样例：
    {'id': '01951bd621b4085cb110fb947a3e8365', 'object': 'chat.completion', 'created': 1739928773, 'model': 'deepseek-ai/DeepSeek-V3', 'choices': [{'index': 0, 'message': {'role': 'assistant', 'content': '分数：-3-，评价：@用户答案中提到了挤压碰撞、冲击、剪切、甩出切割、电机高空坠落等伤害，但未完全覆盖标准答案中的6个伤害。同时，用户答案未提及绞车刹固有风险主要关注的5个方面。建议用户更全面地回答问题，确保覆盖所有关键点。@'}, 'finish_reason': 'stop'}], 'usage': {'prompt_tokens': 153, 'completion_tokens': 71, 'total_tokens': 224}, 'system_fingerprint': ''}
    """
    print(f"deepseek大模型返回结果：{response['choices'][0]['message']['content']}")
    return response['choices'][0]['message']['content']

#四种保存信息到知识库中
'''
知识库sql如下：
CREATE TABLE `knowledgebase` (
  `id` int NOT NULL AUTO_INCREMENT,
  `title` varchar(255) DEFAULT NULL COMMENT '标题',
  `content` longtext COMMENT '正文',
  `label_id` int DEFAULT NULL COMMENT '绑定标签id',
  `user_id` int DEFAULT NULL COMMENT '绑定user的id',
  `type_id` int DEFAULT NULL COMMENT '绑定收藏类型id（是文献、专利、网络还是自定义）',
  `mark_info` varchar(255) DEFAULT NULL COMMENT '对应格式的标注信息，比如文献的doi用来判断是否收藏',
  PRIMARY KEY (`id`)
) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4 COLLATE=utf8mb4_0900_ai_ci;

'''

def add_article_to_knowledge(article_data: dict,label_id:int,user_id:int):
    """
    文献收藏接口，将传入的文献数据插入到 SQL 数据库中
    传入article_data字典样式:
    {'abstract': 'In this paper, we...', 'doi': '10.1016/j.cam.2025.116817', 'id': '686e31a4163c01c8500b6b19', 'keywords': ['(3+1)-dimensional generalized Ito equation', 'Dbar-dressing method', 'Bäcklund transformation', 'Soliton solitons', 'Rational solutions'], 'keywords_zh': None, 'title': 'The ∂̄-Dressing Method, Bäcklund Transformation n and Exact Solutions for a (3+1)-Dimensional Generalized Ito Equation', 'year': 2026, 'is_collected': 1}

    返回:
        bool: 插入成功返回 True，失败返回 False
    """

    # 从字典中提取必要的信息
    title = str(article_data.get('title_zh', article_data['title'])).replace("'", "").replace("\n","")
    abstract = article_data.get('abstract_zh', article_data.get('abstract', ''))
    keywords = article_data.get('keywords_zh',article_data['keywords'])

    #将上述关键信息组成一整段话
    content = '''
    文献信息如下:
    标题：{title}
    摘要：{abstract}
    关键词：{keywords}
    '''.format(title=title,abstract=abstract,keywords=keywords).replace("'", "").replace("\n","")

    # 对于1文献收藏,标注信息为doi
    mark_info = article_data.get('doi')
    sql_sentence =  f"""
    INSERT INTO `knowledgebase` (`title`, `content`, `label_id`, `user_id`, `type_id`, `mark_info`)
    VALUES ('{title}', '{content}', {label_id}, {user_id}, 1, '{mark_info}');
    """
    # print(sql_sentence)
    kyzs_sql.mysql_exec(sql_sentence)

    return {'code': 200, 'msg': 'success', 'data': None}


def add_patent_to_knowledge(patent_data: dict,label_id:int,user_id:int):
    """
    收藏专利
    传入patent_data格式如下：

    {'abstract': {'en': ['The invention discloses.'], 'zh': ['本发明公开了一种']}, 'app_date': '2017-12-08T08:00:00', 'app_kind': 'A', 'app_num': '201711296755', 'applicant': [{'name': '青海煤炭地质一三二勘探队', 'raw_address_info': '810000 青海省西宁市城西区砖厂路13 号', 'sequence': 1}], 'assignee': [{'name': 'QINGHAI COAL GEOLOGY NO 132 ROSPECTING (QING-Non-standard)', 'sequence': 1}], 'auth_date': {'seconds': 1526601600}, 'auth_num': '108049469', 'claims': {'zh': ['1.便携水驱动式管道疏通机，其特征在于，包括冲击室；所述冲击室的内部转动设置有转轴，转轴上固定设置 有叶轮，且转轴的另一端贯穿冲击室向外延伸，转轴的延伸端固定设置有破碎刀；所述冲击室上还连通设置有进水接头，且进水接头的出口与叶轮对应；所述冲击室侧壁开 设有倾斜的通孔，且通孔的进水端高于通孔的出水端。', '2.根据权利要求1所述的便携水驱动式管道疏通机，其特征在于，所述通孔至少为两个，且多个通孔沿冲击室的外壁圆周周向均匀间隔排布。']}, 'country': 'cn', 'cpc': [{'l1': 'E', 'l2': 'E03', 'l3': 'E03C', 'l4': 'E03C1/302', 'raw': 'E03C   1/302'}, {'l1': 'E', 'l2': 'E03', 'l3': 'E03C', 'l4': 'E03C1/306', 'raw': 'E03C   1/306'}], 'description': {'zh': ['便携水驱动式管道疏通机及疏通方法', '技术领域', ]}, 'ep_family_id': '62123076', 'id': '6336263f667297566ca3c494', 'inventor': [{'name': '杨 辰铭', 'sequence': 1}, {'name': '王国志', 'sequence': 2}], 'ipc': [{'l1': 'E', 'l2': 'E03', 'l3': 'E03C', 'l4': 'E03C001/302'}, {'l1': 'E', 'l2': 'E03', 'l3': 'E03C', 'l4': 'E03C001/306'}], 'ipcr': [{'l1': 'E', 'l2': 'E03', 'l3': 'E03C', 'l4': 'E03C1/302'}, {'l1': 'E', 'l2': 'E03', 'l3': 'E03C', 'l4': 'E03C1/306'}], 'priority': [{'country': 'cn', 'date': {'seconds': 1512662400}, 'num': '201711296755'}], 'pub_date': '2018-05-18T08:00:00', 'pub_kind': 'A', 'pub_num': '108049469', 'title': {'en': ['Portable water driving type pipeline dredging machine and dredging method'], 'zh': ['便携水驱动式管道疏通机及疏通方法']}, 'is_collected': 1}


    :param patent_dict: 专利字典
    :param label_id: 标签id
    :param user_id: 用户id
    :return: 收藏结果
    """
    # print(patent_data)
    # 从字典中提取必要的信息
    id = patent_data.get('id')
    title = str(patent_data.get('title')).replace("'", "\"")
    abstract = str(patent_data.get('abstract')).replace("'", "\"")
    claims = str(patent_data.get('claims')).replace("'", "\"")
    country = str(patent_data.get('country')).replace("'", "")
    app_num = str(patent_data.get('app_num')).replace("'", "")
    #app_date 格式为：2021-11-22T08:00:00 ，转换为datetime格式
    app_date = datetime.datetime.strptime(str(patent_data.get('app_date')), '%Y-%m-%dT%H:%M:%S')       
    pub_num = str(patent_data.get('pub_num')).replace("'", "")
    #pub_date 格式为：2021-11-22T08:00:00 ，转换为datetime格式
    pub_date = datetime.datetime.strptime(str(patent_data.get('pub_date')), '%Y-%m-%dT%H:%M:%S')
    pub_kind = str(patent_data.get('pub_kind')).replace("'", "")
    applicant = str(patent_data.get('applicant')).replace("'", "\"")
    assignee = str(patent_data.get('assignee')).replace("'", "\"")
    inventor = str(patent_data.get('inventor')).replace("'", "\"")
    ipcr = str(patent_data.get('ipcr')).replace("'", "\"")       

    #编写有结构的收藏content

    content = f'''
    专利摘要：{abstract}
    专利国家：{country}
    专利申请号：{app_num}
    专利申请日期：{app_date}
    专利公开号：{pub_num}
    专利公开日期：{pub_date}
    专利公开类型：{pub_kind}
    专利申请人：{applicant}
    '''
    #将数据插入到collected_patents专利收藏表中
    sql=f"INSERT INTO knowledgebase (title, content, label_id, user_id, type_id,mark_info) VALUES ('{title}', '{content}', '{label_id}', '{user_id}', 2,'{id}')"
    # print("\033[31m" + sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def add_online_infomation_to_knowledge(online_infomation:dict,label_id:int,user_id:int):
    """
    收藏互联网信息
    :param online_infomation: 互联网信息
    :param label_id: 标签id
    :param user_id: 用户id

    :return: 收藏结果
    """
    # print(online_infomation)
    content= f'''
    日期：{str(online_infomation['date'])}
    标题：{str(online_infomation['title'])}
    内容：{str(online_infomation['content'])}
    来源：{str(online_infomation['source'])}
    '''
    sql=f"INSERT INTO knowledgebase (title, content, label_id, user_id, type_id,mark_info) VALUES ('{online_infomation['title']}', '{content}', '{label_id}', '{user_id}', 3,'网络信息收藏')"
    print("\033[31m" + "收藏互联网信息："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def add_mycontent_to_knowledge(content:str,title:str,label_id:int,user_id:int):
    """
    添加用户自定义内容
    :param content: 用户自定义内容
    :param title: 标题
    :param label_id: 标签id
    :param user_id: 用户id
    :return: 添加结果
    """
    sql=f"INSERT INTO knowledgebase (title,content,label_id,user_id,type_id,mark_info) VALUES ('{title}','{content}',{label_id},{user_id},4,'用户自定义上传信息')"
    print("\033[31m" + "添加用户自定义内容："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def add_mycontent_file_to_knowledge(file_base64_string:str,file_extension:str,title:str,label_id:int,user_id:int):
    """
    添加用户自定义文件内容，读取文件base64，调用方法转为字符串，执行数据库插入
    :param file_base64_string: 文件base64编码字符串
    :param file_extension: 文件扩展名，如pdf、docx、txt等
    :param title: 标题（保留参数，实际使用解析后的文件名）
    :param label_id: 标签id
    :param user_id: 用户id
    :return: 添加结果
    """
    def text_extraction(b64_string, file_type):
        print(f'文件内容:{b64_string[:100]},文件类型:{file_type}')       
        import base64
        import os
        try:
            binary = base64.b64decode(b64_string)
        except:
            return "Invalid Base64 string", ""
        
        # 生成文件名（使用当前时间戳和文件类型）
        import time
        file_name = f"{int(time.time())}.{file_type}"
        
        def pdf_handler(file_binary):
            from pypdf import PdfReader
            with open("temp.pdf", "wb") as f:
                f.write(file_binary)
            reader = PdfReader("temp.pdf")
            full_text = ""
            for page in reader.pages:
                full_text += page.extract_text()
            os.remove("temp.pdf")
            return full_text, f"{int(time.time())}.pdf"
        
        def docx_handler(file_binary):
            from docx import Document
            with open("temp.docx", "wb") as f:
                f.write(file_binary)
            document = Document('temp.docx')
            full_text = ""
            for paragraph in document.paragraphs:
                full_text += paragraph.text
            os.remove("temp.docx")
            return full_text, f"{int(time.time())}.docx"
                    
        def txt_handler(file_binary):
            return file_binary.decode('utf-8'), f"{int(time.time())}.txt"
        
        def ppt_handler(file_binary):
            from pptx import Presentation
            with open("temp.pptx", "wb") as f:
                f.write(file_binary)
            prs = Presentation('temp.pptx')
            full_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
            os.remove("temp.pptx")
            return full_text, f"{int(time.time())}.pptx"
            
        content_string = ""
        match file_type:
            case "txt":
                content_string, file_name = txt_handler(binary)
            case "docx":
                content_string, file_name = docx_handler(binary)
            # case "doc":
            #     content_string, file_name = docx_handler(binary)
            case "pdf":
                content_string, file_name = pdf_handler(binary)
            case "pptx":
                content_string, file_name = ppt_handler(binary)
            case _:
                file_name = f"{int(time.time())}.{file_type}"
 
        return content_string, file_name
    

    print(f'传入base64{file_base64_string[:100]}')
    #先判断这个文件名是否在当前用户数据库中存在了，如果存在则不执行插入操作
    sql=f"SELECT * FROM knowledgebase WHERE title = '{title}' and user_id = {user_id}"
    if kyzs_sql.mysql_exec(sql):
        return {"code": 400, "msg": "数据库当前用户中已存在相同文件名，不执行插入操作", "data": None}
    content_string, file_name = text_extraction(file_base64_string, str(file_extension).replace(".",""))
    # 替换单引号为双引号，避免SQL注入
    content_string = content_string.replace("'", "''")
    # 删除\n，部分pdf会包含\n，
    content_string = content_string.replace("\n", "")
    # 使用text_extraction返回的文件名作为标题
    real_title = title
    sql=f"INSERT INTO knowledgebase (title,content,label_id,user_id,type_id,mark_info) VALUES ('{real_title}','{content_string}',{label_id},{user_id},5,'用户自定义上传文件转换文本内容')"
    print("\033[31m" + "添加用户自定义文件内容："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

    
#更改综述方法，先获取原始文本，替换对应字段，再更新回去
def modify_summary_api(review_id:int, start_position:int, end_position:int, replaced_text:str):
    #获取原始文本
    sql_sentence = f"""
        SELECT review_body FROM review_records
        WHERE id = {review_id}
    """
    original_text = kyzs_sql.mysql_exec(sql_sentence)[0]['review_body']

    #替换文本
    new_text = original_text[:start_position] + replaced_text + original_text[end_position:]
    # print(new_text)

    #更新数据库
    sql_sentence = f"""
        UPDATE review_records
        SET review_body = '{new_text}'
        WHERE id = {review_id}
    """
    kyzs_sql.mysql_exec(sql_sentence)
    return 1

# 全文更改综述方法
def modify_review_new_api(review_id:int, review_body:str):
    #更新数据库
    sql_sentence = f"""
        UPDATE review_records
        SET review_body = '{review_body}'
        WHERE id = {review_id}
    """
    kyzs_sql.mysql_exec(sql_sentence)
    return 1


def delete_summary(id:int):
    sql_sentence = f"""
        DELETE FROM review_records
        WHERE id = {id}
    """
    kyzs_sql.mysql_exec(sql_sentence)
    return 1








def add_format_api(format:dict):
    """
    添加格式要求
    :param format: 格式要求
    :return: 添加结果
    """
    sql=f"INSERT INTO format (class_name, class_text) VALUES ('{format['class_name']}', '{format['class_text']}')"
    print("\033[31m" + "添加格式要求："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def delete_format_api(id:int):
    """
    删除格式要求
    :param id: 格式要求ID
    :return: 删除结果
    """
    sql=f"DELETE FROM format WHERE id={id}"
    print("\033[31m" + "删除格式要求："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}


def generate_word_api(id:int):
    """
    传入综述id，生成Word文档
    :param id: 综述id
    :return: 生成的word的二进制文本
    """

    #从数据库中获取综述内容
    sql=f"SELECT * FROM review_records WHERE id={id}"
    print("\033[31m" + "获取综述内容用于生成word："+sql + "\033[0m")
    data = kyzs_sql.mysql_exec(sql)
    if data:
        md_text = data[0]['review_body']
    else:
        return {'code': 400, 'msg': '未找到该综述', 'data': None}

    # 先将文本转为.md文件并保存为临时文件
    md_text = md_text.replace('\\n', '\n')

    # 删除文本中的```markdown和```
    md_text = md_text.replace('```markdown', '')
    md_text = md_text.replace('```', '')
    
    md_file = 'temp.md'
    with open(md_file, 'w', encoding='utf-8') as f:
        f.write(md_text)
    

    # 转换为 Word 文档
    output_path = "temp_output_word.docx"
    pypandoc.convert_file(md_file, to='docx', outputfile=output_path)
    # 读取生成的 Word 文档
    with open(output_path, 'rb') as f:
        word_binary = f.read()

    # 清理临时文件
    os.remove(md_file)
    os.remove(output_path)

    #将编码转为base64
    import base64
    word_64 = base64.b64encode(word_binary).decode('utf-8')
    print('生成word文档成功')
    return {'code': 200, 'msg': 'success', 'data': word_64}




def generate_fuwenben_word_api(id: int):
    """
    传入综述id，生成富文本Word文档
    :param id: 综述id
    :return: 生成的word的二进制文本（base64编码字符串）
    """
    # 从数据库中获取综述内容
    sql = f"SELECT * FROM review_records WHERE id={id}"
    print("\033[31m" + "获取综述内容用于生成word：" + sql + "\033[0m")
    data = kyzs_sql.mysql_exec(sql)
    if not data:
        return {'code': 400, 'msg': '未找到该综述', 'data': None}
    
    html_text = data[0]['review_body']  # 这里的 review_body 是 HTML 富文本字符串

    # 使用 BeautifulSoup 解析 HTML 并转换为 python-docx 文档
    from docx import Document
    from bs4 import BeautifulSoup
    from io import BytesIO
    import base64

    doc = Document()
    
    soup = BeautifulSoup(html_text, 'html.parser')
    
    def process_element(element, parent_paragraph=None):
        """递归处理 HTML 元素，添加到 Word 文档中"""
        if element.name is None:  # 文本节点
            text = str(element)
            if parent_paragraph:
                run = parent_paragraph.add_run(text)
                return run
            return None
        
        if element.name in ['p', 'div', 'br']:  # 段落或换行
            if parent_paragraph and parent_paragraph.runs:
                # 当前段落已有内容，先结束它
                pass
            new_p = doc.add_paragraph()
            for child in element.children:
                process_element(child, new_p)
            return None
        
        elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            level = int(element.name[1])  # h1 -> 1, h2 -> 2 等
            heading = doc.add_heading(element.get_text(strip=False), level=level)
            return None
        
        elif element.name == 'strong' or element.name == 'b':
            run = None
            if parent_paragraph:
                run = parent_paragraph.add_run(element.get_text())
                run.bold = True
            for child in element.children:
                child_run = process_element(child, parent_paragraph)
                if child_run:
                    child_run.bold = True
            return run
        
        elif element.name == 'em' or element.name == 'i':
            run = None
            if parent_paragraph:
                run = parent_paragraph.add_run(element.get_text())
                run.italic = True
            for child in element.children:
                child_run = process_element(child, parent_paragraph)
                if child_run:
                    child_run.italic = True
            return run
        
        else:  # 其他标签（如 span、ul 等），这里简单当作普通文本处理
            for child in element.children:
                process_element(child, parent_paragraph)
            return None

    # 从根开始处理
    for child in soup.children:
        process_element(child)

    # 将 docx 保存到内存中
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    
    # 转为 base64
    word_b64 = base64.b64encode(buffer.read()).decode('utf-8')
    
    return {'code': 200, 'msg': 'success', 'data': word_b64}
# 个人自定义资料

def add_mycontent_with_ai_api(content:str):
    """
    添加用户自定义内容，且通过大模型做总结后再保存
    :param content: 用户自定义内容
    :return: 添加结果
    """

    #通过大模型做总结
    summary = siliconflow_deepseek_answer("请对以下文本做总结："+content)
    return add_mycontent_api(summary)
def get_all_mycontent_api():
    """
    获取所有用户自定义内容
    :return: 用户自定义内容列表
    """
    sql=f"SELECT * FROM collected_mycontents"
    print("\033[31m" + "获取所有用户自定义内容："+sql + "\033[0m")
    data = kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': data}

def delete_mycontent_api(id:int):
    """
    删除用户自定义内容
    :param id: 用户自定义内容ID
    :return: 删除结果
    """
    sql=f"DELETE FROM collected_mycontents WHERE id={id}"
    print("\033[31m" + "删除用户自定义内容："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}



if __name__ == "__main__":
    import os
    # 切换目录到c:\Users\Administrator\Desktop\科研情报-临时部署\
    os.chdir("c:\\Users\\Administrator\\Desktop\\科研情报-临时部署\\")
    print(f"当前工作目录: {os.getcwd()}")
    
    # 跳过验证阶段
    print(get_online_infomation_api("中国石化是什么"))