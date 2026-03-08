import requests
from app.core.database import db as kyzs_sql
from app.config import settings
import json
import datetime
import os
import pypandoc
import base64
import hashlib
import time
import urllib.parse
import xml.etree.ElementTree as ET


def get_articles_from_oillink(keywords_list: list, page: int, size: int, user_id: int):
    """
    从 OilLink 检索文献，并标记每篇文章是否已被当前用户收藏（is_collected）。
    用 DOI 作为唯一标识匹配 knowledgebase 表。
    """
    url = "http://data.oillink.com/api/shengli/articlesearch/index"
    params = {"keywords": f'{keywords_list}', "page": page, "size": size}
    try:
        response = requests.get(url, params=params)
        response.raise_for_status()
        data = response.json()['data']
        if data is None:
            return []
        for article in data:
            doi = article.get('doi')
            if not doi:
                article['is_collected'] = 0
                continue
            result = kyzs_sql.mysql_exec(
                "SELECT COUNT(*) as cnt FROM knowledgebase WHERE mark_info=%s AND type_id=1 AND user_id=%s",
                (doi, user_id)
            )
            article['is_collected'] = 1 if result and result[0]['cnt'] > 0 else 0
        return data
    except requests.RequestException as e:
        print(f"请求出错: {e}")
        return None


def get_patents_from_oillink(query: str, page: int = 0, size: int = 5):
    url = "http://data.oillink.com/api/shengli/patentsearch/index"
    params = {"query": query, "page": page, "size": size}
    try:
        response = requests.post(url, params=params)
        response.raise_for_status()
        data = response.json()
        if data['code'] != 200:
            return []
        data = data['data']
        new_data = []
        for item in data:
            i = item[0]
            patent_id = i['id']
            id_result = kyzs_sql.mysql_exec(
                "SELECT COUNT(*) as cnt FROM knowledgebase WHERE mark_info=%s AND type_id=2",
                (patent_id,)
            )
            i['is_collected'] = 1 if id_result and id_result[0]['cnt'] else 0
            i['app_date'] = datetime.datetime.fromtimestamp(i['app_date']['seconds'])
            i['pub_date'] = datetime.datetime.fromtimestamp(i['pub_date']['seconds'])
            new_data.append(i)
        return new_data
    except requests.RequestException as e:
        print(f"请求出错: {e}")
        return None


def get_online_infomation_api(keyword: str):
    def xunfei_online_search(question):
        url = "https://spark-api-open.xf-yun.com/v1/chat/completions"
        print(f'方法内：讯飞网络知识检索关键词：{question}', flush=True)
        tools = [
            {
                "type": "function",
                "function": {
                    "name": "get_research_information",
                    "description": "获取互联网科研信息",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "信息标题"},
                            "content": {"type": "string", "description": "信息内容，约1000字，倾向于真实案例和数据"},
                            "source": {"type": "string", "description": "信息来源"},
                            "date": {"type": "string", "description": "信息日期，格式为YYYY-MM-DD"},
                            "author": {"type": "string", "description": "信息作者"}
                        },
                        "required": ["title", "content", "source", "date", "author"]
                    }
                }
            },
            {
                "type": "web_search",
                "web_search": {"enable": True, "show_ref_label": True, "search_mode": "deep"}
            }
        ]
        data = {
            "model": "4.0Ultra",
            "user": "default_user",
            "messages": [
                {"role": "system", "content": "你是知识渊博的助理，能够获取互联网科研信息"},
                {"role": "user", "content": question}
            ],
            "temperature": 1,
            "top_k": 6,
            "stream": False,
            "max_tokens": 5000,
            "tools": tools,
            "tool_choice": {"type": "function", "function": {"name": "get_research_information"}},
        }
        header = {"Authorization": f"Bearer {settings.xunfei_api_key}"}
        print(f'准备执行请求：{url}，请求头：{header}，请求数据：{data}', flush=True)
        try:
            response = requests.post(url, headers=header, json=data)
            print(f'请求已发送，等待响应...', flush=True)
            return response.json()
        except Exception as e:
            print(f'发送请求时出错：{e}', flush=True)
            return None

    return xunfei_online_search(keyword)


def get_article_from_juhe_api(keywords: list, date: str = None, page=1, size=10, sort=3):
    def md5_encrypt(text: str) -> str:
        md5 = hashlib.md5()
        md5.update(text.encode('utf-8'))
        return md5.hexdigest()

    exps = [f"关键词:{k}" for k in keywords]
    if date:
        exps.append(f"年份:{date}")
    exp = urllib.parse.quote(" AND ".join(exps))
    params = "{}|{}|{}|{}".format(exp, page, size, sort)
    print('表达式：', params)
    times = "{}".format(int(time.time() * 1000))
    en_params = base64.b64encode(params.encode('utf-8')).decode('utf-8')
    sign = md5_encrypt("{}|{}|{}".format(en_params, settings.juhe_token, times))
    headers = {"User-Agent": "test", "Accept-Encoding": "gzip, deflate, br", "Connection": "keep-alive"}
    query_params = {"params": en_params, "sign": sign, "token": settings.juhe_token, "times": times}
    try:
        with requests.Session() as session:
            url = "http://61.128.134.70:6655/groups/ask/search"
            session.get(url, headers=headers, params=query_params, timeout=10)
            response = session.get(url, headers=headers, params=query_params, timeout=10)
            print("Raw Response:")
            print(response.text)
            response.raise_for_status()
            result = response.json()['result']
            return [
                {
                    "标题": r['title'], "关键词": r['keyword'], "年份": r['year'],
                    "摘要": r['content'], "DOI": r['doi'],
                    "下载链接": f"http://61.128.134.70:6655{r['full_source']}"
                }
                for r in result['records']
            ]
    except Exception as e:
        print(f"重庆聚合请求出错: {e}")
        return None


def get_article_from_wanfang_api(Datewithin=None, Keywords=[], StartRecord=1, MaximumRecords=10):
    exps = list(Keywords)
    if Datewithin:
        exps.append(f'Date within "{Datewithin}-01-01 {Datewithin}-12-31"')
    exp = " and ".join(exps)
    print(f"检索表达式：{exp}")

    def send_soap_request_paper(exp, startRecord=1, maximumRecords=10):
        url = "http://10.68.16.2/S/SRW/Paper.asmx"
        headers = {"Host": "10.68.16.2", "Content-Type": "text/xml; charset=utf-8", "SOAPAction": "searchRetrieve"}
        soap_body = f"""<?xml version="1.0" encoding="utf-8"?>
            <soap:Envelope xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"
                        xmlns:xsd="http://www.w3.org/2001/XMLSchema"
                        xmlns:soap="http://schemas.xmlsoap.org/soap/envelope/">
            <soap:Body>
                <searchRetrieveRequest xmlns="http://www.loc.gov/zing/srw/">
                    <version>1.2</version>
                    <query>{exp}</query>
                    <operation>searchretrieve</operation>
                    <startRecord>{startRecord}</startRecord>
                    <maximumRecords>{maximumRecords}</maximumRecords>
                </searchRetrieveRequest>
            </soap:Body>
            </soap:Envelope>"""
        try:
            response = requests.post(url, data=soap_body, headers=headers)
            response.raise_for_status()
            print("Response Status Code:", response.status_code)
            return response.text
        except requests.exceptions.RequestException as e:
            print("Error sending SOAP request:", e)
            return None

    response = send_soap_request_paper(exp, startRecord=StartRecord, maximumRecords=MaximumRecords)
    with open("raw_response.txt", "w") as f:
        f.write(f"{response}")

    def parse_academic_papers(soap_response):
        namespaces = {
            "soap": "http://schemas.xmlsoap.org/soap/envelope/",
            "srw": "http://www.loc.gov/zing/srw/",
            "dc": "info:srw/schema/1/dc-v1.1",
            "srw_dc": "info:srw/schema/1/dc-v1.1"
        }
        print(soap_response)
        root = ET.fromstring(soap_response)
        number_element = root.find('.//srw:numberOfRecords', namespaces=namespaces)
        if number_element is not None:
            number_of_records = number_element.text
            print(f"共解析到论文数量: {number_of_records}")
        else:
            print("未找到numberOfRecords元素")
            number_of_records = "0"
        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record", namespaces=namespaces
        )
        papers = []
        for record in record_nodes:
            dc_node = record.find("srw:recordData/srw_dc:dc", namespaces=namespaces)
            def _get(tag):
                el = dc_node.find(tag, namespaces)
                return el.text.strip() if el is not None and el.text else ""
            info = {
                "标题": _get("dc:title"), "关键词": _get("dc:Subject"),
                "摘要": _get("dc:Description"), "发表时间": _get("dc:Date"),
                "DOI": _get("dc:Identifier"),
            }
            info["发表时间"] = info["发表时间"].split("-")[0]
            papers.append(info)
        return (papers, number_of_records)

    return parse_academic_papers(response)


def wangfang_patent(Datewithin=None, patent_name=[], StartRecord=1, MaximumRecords=10):
    exps = list(patent_name)
    if Datewithin:
        exps.append(f'F_PublicationDate within "{Datewithin}-01-01 {Datewithin}-12-31"')
    exp = " and ".join(exps)
    # print(f"检索表达式：{exp}")  # 原版注释掉了，保持一致

    def send_soap_request_patent(exp, startRecord=1, maximumRecords=10):
        url = "http://10.68.16.2/S/SRW/Patent.asmx"
        headers = {"Host": "10.68.16.2", "Content-Type": "text/xml; charset=utf-8", "SOAPAction": "searchRetrieve"}
        soap_body = f"""<?xml version="1.0" encoding="utf-8"?>
            <soap:Envelope xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"
                        xmlns:xsd="http://www.w3.org/2001/XMLSchema"
                        xmlns:soap="http://schemas.xmlsoap.org/soap/envelope/">
            <soap:Body>
                <searchRetrieveRequest xmlns="http://www.loc.gov/zing/srw/">
                    <version>1.2</version>
                    <query>{exp}</query>
                    <operation>searchretrieve</operation>
                    <startRecord>{startRecord}</startRecord>
                    <maximumRecords>{maximumRecords}</maximumRecords>
                </searchRetrieveRequest>
            </soap:Body>
            </soap:Envelope>"""
        try:
            response = requests.post(url, data=soap_body, headers=headers)
            response.raise_for_status()
            print("Response Status Code:", response.status_code)
            return response.text
        except requests.exceptions.RequestException as e:
            print("Error sending SOAP request:", e)
            return None

    response = send_soap_request_patent(exp, startRecord=StartRecord, maximumRecords=MaximumRecords)

    def parse_patents(soap_response) -> tuple:
        namespaces = {
            "soap": "http://schemas.xmlsoap.org/soap/envelope/",
            "srw": "http://www.loc.gov/zing/srw/",
            "dc": "info:srw/schema/1/dc-v1.1",
            "srw_dc": "info:srw/schema/1/dc-v1.1"
        }
        # print(soap_response)  # 原版注释掉了，保持一致
        root = ET.fromstring(soap_response)
        number_element = root.find('.//srw:numberOfRecords', namespaces=namespaces)
        if number_element is not None:
            number_of_records = number_element.text
            print(f"共解析到论文数量: {number_of_records}")
        else:
            print("未找到numberOfRecords元素")
            number_of_records = "0"
        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record", namespaces=namespaces
        )
        papers = []
        for record in record_nodes:
            dc_node = record.find("srw:recordData/srw_dc:dc", namespaces=namespaces)
            def _get(tag):
                el = dc_node.find(tag, namespaces)
                return el.text.strip() if el is not None and el.text else ""
            papers.append({
                "申请号": _get("dc:ApplicationNo"), "申请日": _get("dc:ApplicationDate"),
                "公开号": _get("dc:PublicationNo"), "公开日": _get("dc:PublicationDate"),
                "专利名称": _get("dc:PatentName"), "申请人": _get("dc:Applicant"),
                "发明人": _get("dc:Inventor"), "IPC分类号": _get("dc:ClassMain"),
                "摘要": _get("dc:Abstract"), "权利要求": _get("dc:SignoryItem"),
            })
        return (papers, number_of_records)

    return parse_patents(response)


def siliconflow_deepseek_answer(question):
    """
    调用 SiliconFlow 托管的 DeepSeek-V3 模型进行问答。
    stream=False 为同步调用，适合需要完整响应的场景（翻译、综述生成等）。
    流式对话场景请使用 AnythingLLM 的 stream-chat 接口。
    """
    url = "https://api.siliconflow.cn/v1/chat/completions"
    payload = {
        "model": "Pro/deepseek-ai/DeepSeek-V3",
        "messages": [{"role": "user", "content": str(question)}],
        "stream": False,
        "max_tokens": 8192,
        "stop": ["null"],
        "temperature": 0.7,
        "top_p": 0.7,
        "top_k": 50,
        "frequency_penalty": 0.5,
        "n": 1,
        "response_format": {"type": "text"},
        "tools": [
            {
                "type": "function",
                "function": {
                    "description": "<string>",
                    "name": "<string>",
                    "parameters": {},
                    "strict": False
                }
            }
        ]
    }
    headers = {
        "Authorization": f"Bearer {settings.siliconflow_api_key}",
        "Content-Type": "application/json"
    }
    response = requests.post(url, json=payload, headers=headers)
    content = response.json()['choices'][0]['message']['content']
    print(f"deepseek返回: {content}")
    return content


def add_article_to_knowledge(article_data: dict, label_id: int, user_id: int):
    title = str(article_data.get('title_zh', article_data['title'])).replace("\n", "")
    abstract = article_data.get('abstract_zh', article_data.get('abstract', ''))
    keywords = article_data.get('keywords_zh', article_data['keywords'])
    content = f"文献信息如下:    标题：{title}    摘要：{abstract}    关键词：{keywords}    "
    mark_info = article_data.get('doi')
    kyzs_sql.mysql_exec(
        "INSERT INTO `knowledgebase` (title, content, label_id, user_id, type_id, mark_info) VALUES (%s,%s,%s,%s,1,%s)",
        (title, content, label_id, user_id, mark_info)
    )
    return {'code': 200, 'msg': 'success', 'data': None}


def add_patent_to_knowledge(patent_data: dict, label_id: int, user_id: int):
    id = patent_data.get('id')
    title = str(patent_data.get('title'))
    abstract = str(patent_data.get('abstract'))
    country = str(patent_data.get('country'))
    app_num = str(patent_data.get('app_num'))
    app_date = datetime.datetime.strptime(str(patent_data.get('app_date')), '%Y-%m-%dT%H:%M:%S')
    pub_num = str(patent_data.get('pub_num'))
    pub_date = datetime.datetime.strptime(str(patent_data.get('pub_date')), '%Y-%m-%dT%H:%M:%S')
    pub_kind = str(patent_data.get('pub_kind'))
    applicant = str(patent_data.get('applicant'))
    content = (
        f"专利摘要：{abstract}\n专利国家：{country}\n专利申请号：{app_num}\n"
        f"专利申请日期：{app_date}\n专利公开号：{pub_num}\n专利公开日期：{pub_date}\n"
        f"专利公开类型：{pub_kind}\n专利申请人：{applicant}\n"
    )
    kyzs_sql.mysql_exec(
        "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (%s,%s,%s,%s,2,%s)",
        (title, content, label_id, user_id, id)
    )
    return {'code': 200, 'msg': 'success', 'data': None}


def add_online_infomation_to_knowledge(online_infomation: dict, label_id: int, user_id: int):
    content = (
        f"日期：{online_infomation['date']}\n标题：{online_infomation['title']}\n"
        f"内容：{online_infomation['content']}\n来源：{online_infomation['source']}\n"
    )
    sql = "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (%s,%s,%s,%s,3,'网络信息收藏')"
    print("\033[31m" + "收藏互联网信息：" + sql + "\033[0m")
    kyzs_sql.mysql_exec(sql, (online_infomation['title'], content, label_id, user_id))
    return {'code': 200, 'msg': 'success', 'data': None}


def add_mycontent_to_knowledge(content: str, title: str, label_id: int, user_id: int):
    sql = "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (%s,%s,%s,%s,4,'用户自定义上传信息')"
    print("\033[31m" + "添加用户自定义内容：" + sql + "\033[0m")
    kyzs_sql.mysql_exec(sql, (title, content, label_id, user_id))
    return {'code': 200, 'msg': 'success', 'data': None}


def add_mycontent_file_to_knowledge(file_base64_string: str, file_extension: str, title: str, label_id: int, user_id: int):
    """
    解析上传文件（PDF/DOCX/TXT/PPTX）并将文本内容存入知识库。
    文件以 base64 字符串传输，解码后写临时文件处理，处理完即删除。
    同一用户下文件标题唯一，重复上传会返回 400。
    """
    def text_extraction(b64_string, file_type):
        print(f'文件内容:{b64_string[:100]},文件类型:{file_type}')
        try:
            binary = base64.b64decode(b64_string)
        except Exception:
            return "Invalid Base64 string", ""

        def pdf_handler(file_binary):
            from pypdf import PdfReader
            with open("temp.pdf", "wb") as f:
                f.write(file_binary)
            reader = PdfReader("temp.pdf")
            full_text = "".join(page.extract_text() or "" for page in reader.pages)
            os.remove("temp.pdf")
            return full_text, f"{int(time.time())}.pdf"

        def docx_handler(file_binary):
            from docx import Document
            with open("temp.docx", "wb") as f:
                f.write(file_binary)
            document = Document('temp.docx')
            full_text = "".join(p.text for p in document.paragraphs)
            os.remove("temp.docx")
            return full_text, f"{int(time.time())}.docx"

        def txt_handler(file_binary):
            return file_binary.decode('utf-8'), f"{int(time.time())}.txt"

        def ppt_handler(file_binary):
            from pptx import Presentation
            with open("temp.pptx", "wb") as f:
                f.write(file_binary)
            prs = Presentation('temp.pptx')
            full_text = "".join(
                shape.text + "\n"
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            os.remove("temp.pptx")
            return full_text, f"{int(time.time())}.pptx"

        handlers = {"txt": txt_handler, "docx": docx_handler, "pdf": pdf_handler, "pptx": ppt_handler}
        handler = handlers.get(file_type)
        if handler:
            return handler(binary)
        return "", f"{int(time.time())}.{file_type}"

    print(f'传入base64{file_base64_string[:100]}')
    existing = kyzs_sql.mysql_exec(
        "SELECT id FROM knowledgebase WHERE title=%s AND user_id=%s", (title, user_id)
    )
    if existing:
        return {"code": 400, "msg": "数据库当前用户中已存在相同文件名，不执行插入操作", "data": None}

    content_string, _ = text_extraction(file_base64_string, str(file_extension).replace(".", ""))
    # 删除\n，部分pdf会包含\n
    content_string = content_string.replace("\n", "")
    kyzs_sql.mysql_exec(
        "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (%s,%s,%s,%s,5,'用户自定义上传文件转换文本内容')",
        (title, content_string, label_id, user_id)
    )
    return {'code': 200, 'msg': 'success', 'data': None}


def modify_summary_api(review_id: int, start_position: int, end_position: int, replaced_text: str):
    result = kyzs_sql.mysql_exec(
        "SELECT review_body FROM review_records WHERE id=%s", (review_id,)
    )
    original_text = result[0]['review_body']
    new_text = original_text[:start_position] + replaced_text + original_text[end_position:]
    kyzs_sql.mysql_exec(
        "UPDATE review_records SET review_body=%s WHERE id=%s", (new_text, review_id)
    )
    return 1


def modify_review_new_api(review_id: int, review_body: str):
    kyzs_sql.mysql_exec(
        "UPDATE review_records SET review_body=%s WHERE id=%s", (review_body, review_id)
    )
    return 1


def delete_summary(id: int):
    kyzs_sql.mysql_exec("DELETE FROM review_records WHERE id=%s", (id,))
    return 1


def generate_word_api(id: int):
    """
    将综述正文（Markdown 格式）转换为 Word 文档，以 base64 返回。
    依赖系统安装的 pandoc，需提前执行 brew install pandoc（macOS）或对应安装命令。
    """
    sql = "SELECT * FROM review_records WHERE id=%s"
    print("\033[31m" + "获取综述内容用于生成word：" + sql + "\033[0m")
    data = kyzs_sql.mysql_exec(sql, (id,))
    if not data:
        return {'code': 400, 'msg': '未找到该综述', 'data': None}

    md_text = data[0]['review_body']
    md_text = md_text.replace('\\n', '\n').replace('```markdown', '').replace('```', '')

    with open('temp.md', 'w', encoding='utf-8') as f:
        f.write(md_text)

    output_path = "temp_output_word.docx"
    pypandoc.convert_file('temp.md', to='docx', outputfile=output_path)
    with open(output_path, 'rb') as f:
        word_binary = f.read()

    os.remove('temp.md')
    os.remove(output_path)

    word_64 = base64.b64encode(word_binary).decode('utf-8')
    print('生成word文档成功')
    return {'code': 200, 'msg': 'success', 'data': word_64}


def generate_fuwenben_word_api(id: int):
    sql = "SELECT * FROM review_records WHERE id=%s"
    print("\033[31m" + "获取综述内容用于生成word：" + sql + "\033[0m")
    data = kyzs_sql.mysql_exec(sql, (id,))
    if not data:
        return {'code': 400, 'msg': '未找到该综述', 'data': None}

    html_text = data[0]['review_body']

    from docx import Document
    from bs4 import BeautifulSoup
    from io import BytesIO

    doc = Document()
    soup = BeautifulSoup(html_text, 'html.parser')

    def process_element(element, parent_paragraph=None):
        if element.name is None:
            if parent_paragraph:
                parent_paragraph.add_run(str(element))
            return
        if element.name in ['p', 'div', 'br']:
            new_p = doc.add_paragraph()
            for child in element.children:
                process_element(child, new_p)
        elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            doc.add_heading(element.get_text(strip=False), level=int(element.name[1]))
        elif element.name in ['strong', 'b']:
            if parent_paragraph:
                run = parent_paragraph.add_run(element.get_text())
                run.bold = True
        elif element.name in ['em', 'i']:
            if parent_paragraph:
                run = parent_paragraph.add_run(element.get_text())
                run.italic = True
        else:
            for child in element.children:
                process_element(child, parent_paragraph)

    for child in soup.children:
        process_element(child)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    word_b64 = base64.b64encode(buffer.read()).decode('utf-8')
    return {'code': 200, 'msg': 'success', 'data': word_b64}
