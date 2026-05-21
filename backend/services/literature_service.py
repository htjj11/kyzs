import requests
import json
import datetime
import os
import base64
import time
import xml.etree.ElementTree as ET
from core.sqlLiteExec import sqlite_execute

#添加文献到知识库
def add_article_to_knowledge(article_data: dict, label_id: int, user_id: int):
    title = str(article_data.get('title_zh', article_data['title'])).replace("\n", "")
    abstract = article_data.get('abstract_zh', article_data.get('abstract', ''))
    keywords = article_data.get('keywords_zh', article_data['keywords'])
    content = f"文献信息如下:    标题：{title}    摘要：{abstract}    关键词：{keywords}    "
    mark_info = article_data.get('doi')
    sqlite_execute(
        "INSERT INTO `knowledgebase` (title, content, label_id, user_id, type_id, mark_info) VALUES (?,?,?,?,?,?)",
        (title, content, label_id, user_id, 1, mark_info)
    )
    return {'code': 200, 'msg': 'success', 'data': None}

#添加专利到知识库
def add_patent_to_knowledge(patent_data: dict, label_id: int, user_id: int):
    id = patent_data.get('id')
    title = str(patent_data.get('title'))
    abstract = str(patent_data.get('abstract'))
    country = str(patent_data.get('country'))
    app_num = str(patent_data.get('app_num'))
    app_date = datetime.datetime.strptime(str(patent_data.get('app_date')), '%Y-%m-%dT%H:%M:%S')
    pub_num = str(patent_data.get('pub_num'))
    pub_date = datetime.datetime.strptime(str(patent_data.get('pub_date')), '%Y-%m-%dT%H:%M:%S')
    pub_kind = str(patent_data.get('pub_kind'))
    applicant = str(patent_data.get('applicant'))
    content = (
        f"专利摘要：{abstract}\n专利国家：{country}\n专利申请号：{app_num}\n"
        f"专利申请日期：{app_date}\n专利公开号：{pub_num}\n专利公开日期：{pub_date}\n"
        f"专利公开类型：{pub_kind}\n专利申请人：{applicant}\n"
    )
    sqlite_execute(
        "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (?,?,?,?,?,?)",
        (title, content, label_id, user_id, 2, id)
    )
    return {'code': 200, 'msg': 'success', 'data': None}

#添加网络信息到知识库
def add_online_infomation_to_knowledge(online_infomation: dict, label_id: int, user_id: int):
    content = (
        f"日期：{online_infomation['date']}\n标题：{online_infomation['title']}\n"
        f"内容：{online_infomation['content']}\n来源：{online_infomation['source']}\n"
    )
    sql = "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (?,?,?,?,3,'网络信息收藏')"
    print("\033[31m" + "收藏互联网信息：" + sql + "\033[0m")
    sqlite_execute(sql, (online_infomation['title'], content, label_id, user_id))
    return {'code': 200, 'msg': 'success', 'data': None}

#添加用户自定义内容到知识库
def add_mycontent_to_knowledge(content: str, title: str, label_id: int, user_id: int):
    sql = "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (?,?,?,?,4,'用户自定义上传信息')"
    print("\033[31m" + "添加用户自定义内容：" + sql + "\033[0m")
    sqlite_execute(sql, (title, content, label_id, user_id))
    return {'code': 200, 'msg': 'success', 'data': None}

#添加用户自定义文件到知识库
def add_mycontent_file_to_knowledge(file_base64_string: str, file_extension: str, title: str, label_id: int, user_id: int):
    """
    解析上传文件列表（PDF/DOCX/TXT/PPTX）并将文本内容存入知识库。
    文件以 base64 字符串传输，保存到file_data目录中，并重命名为纯数字（以毫秒级别时间戳命名）。
    同一用户下文件标题唯一，重复上传会返回 400。
    """
    print(f'传入base64{file_base64_string[:100]}')
    existing = sqlite_execute(
        "SELECT id FROM knowledgebase WHERE title=? AND user_id=?", (title, user_id)
    )
    if existing:
        return {"code": 400, "msg": "数据库当前用户中已存在相同文件名，不执行插入操作", "data": None}

    os.makedirs(os.path.join("file_data", "personal_db"), exist_ok=True)
    file_ext = str(file_extension).replace(".", "")
    new_filename = f"{int(time.time() * 1000)}.{file_ext}"
    rel_filename = f"personal_db/{new_filename}"
    file_path = os.path.join("file_data", rel_filename)

    max_bytes = 200 * 1024 * 1024  # 与前端「资料上传」单文件上限一致
    try:
        binary = base64.b64decode(file_base64_string)
        if len(binary) > max_bytes:
            return {"code": 400, "msg": f"单文件不能超过 {max_bytes // (1024 * 1024)}MB", "data": None}
        with open(file_path, "wb") as f:
            f.write(binary)
    except Exception:
        return {"code": 500, "msg": "Invalid Base64 string", "data": None}

    def text_extraction(path, file_type):
        def pdf_handler(p):
            from pypdf import PdfReader
            reader = PdfReader(p)
            return "".join(page.extract_text() or "" for page in reader.pages)

        def docx_handler(p):
            import zipfile
            import xml.etree.ElementTree as ET
            try:
                with zipfile.ZipFile(p) as docx:
                    xml_content = docx.read('word/document.xml')
                tree = ET.XML(xml_content)
                WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
                texts = [node.text for node in tree.iter(WORD_NAMESPACE + 't') if node.text]
                return "".join(texts).replace('\n', '')
            except Exception:
                from docx import Document
                document = Document(p)
                texts = [p_obj.text for p_obj in document.paragraphs]
                for table in document.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
                return "".join(texts).replace('\n', '')

        def txt_handler(p):
            with open(p, "rb") as f:
                return f.read().decode('utf-8')

        def ppt_handler(p):
            from pptx import Presentation
            prs = Presentation(p)
            return "".join(
                shape.text + "\n"
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

        handlers = {"txt": txt_handler, "docx": docx_handler, "pdf": pdf_handler, "pptx": ppt_handler}
        handler = handlers.get(file_type)
        if handler:
            return handler(path)
        return ""

    content_string = text_extraction(file_path, file_ext)
    # 删除\n，部分pdf会包含\n
    content_string = content_string.replace("\n", "")
    
    mark_info_dict = {
        "filename": rel_filename,
        "original_filename": f"{title}.{file_ext}"
    }
    mark_info_str = json.dumps(mark_info_dict, ensure_ascii=False)
    
    sqlite_execute(
        "INSERT INTO knowledgebase (title, content, label_id, user_id, type_id, mark_info) VALUES (?,?,?,?,5,?)",
        (title, content_string, label_id, user_id, mark_info_str)
    )
    return {'code': 200, 'msg': 'success', 'data': None}

