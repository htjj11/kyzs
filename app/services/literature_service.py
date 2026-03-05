'''
此处定义所有的接口的方法
'''

import requests
from app.core.database import db as kyzs_sql
import json
import datetime
import os
import pypandoc
import base64
import hashlib
import time
import urllib.parse


#三种从oilink获取信息的方法
def get_articles_from_oillink(keywords_list:list,page:int,size:int,user_id:int):
    """
    oillink的文献检索接口
    执行 curl 请求，从 http://data.oillink.com 获取文章搜索结果
    
    返回:
        requests.Response 对象，包含请求的响应
        格式为：{'code': 200, 'msg': 'success', 'data': [{'abstract_zh': '直线型振动筛和...', 'id': '6860e2f5163c01c850d9987f', 'keywords': None, 'keywords_zh': ['钻井液振动筛', '双振型', '离散单元法', 'EDEM', '筛分效率'], 'title_zh': '基于离散单元法的双振型下钻井液振动筛不 同颗粒筛分效率对比研究 后印本', 'year': 2025}, {'abstract': 'The primary ..', 'abstract_zh': '的目的...', 'doi': '10.3969/j.issn.1000-6532.2025.02.012', 'id': '68108681163c01c850d8d66d', 'keywords': ['Mining processing engineering', 'Flotation', 'Barite', 'Waste drilling fluid', 'Inhibitors', 'Collectors', 'Mechanism study'], 'keywords_zh': ['矿物加工工程', '浮选', '重晶石', '废弃钻井液', '抑制剂', '捕收剂', '机理研究'], 'title': 'Flotation Recovery of Barite in Waste Drilling Fluids Using Novel Inhibitors', 'title_zh': '新型抑制剂对废弃钻井液中重晶石的浮选回收', 'year': 2025}]}
    """
    url = "http://data.oillink.com/api/shengli/articlesearch/index"
    params = {
        "keywords": f'{keywords_list}',
        "page": page,
        "size": size
    }
    try:
        response = requests.get(url, params=params)
        response.raise_for_status()

        #处理返回数据总的data列表，添加一个字段is_collected,遍历其中的doi是否在收藏列表中，如果在则为1，否则为0
        data = response.json()['data']

        if data == None:
            return []

        for article in data:
            doi = article.get('doi')
            if not doi:
                continue
            sql_sentence = f"""
            SELECT COUNT(*) FROM knowledgebase WHERE mark_info LIKE '{doi}' AND type_id = 1 AND user_id = {user_id};
            """
            result = kyzs_sql.mysql_exec(sql_sentence)
            if result[0]['COUNT(*)'] > 0:
                article['is_collected'] = 1
            else:
                article['is_collected'] = 0    
        return data


    except requests.RequestException as e:
        print(f"请求出错: {e}")
        return None

def get_patents_from_oillink(query:str, page:int = 0, size:int = 5):
    """
    专利检索执行 POST 请求，从 http://data.oillink.com 获取专利搜索结果
    """
    url = "http://data.oillink.com/api/shengli/patentsearch/index"
    params = {
        "query": query,
        "page": page,
        "size": size
    }
    try:
        response = requests.post(url, params=params)
        response.raise_for_status()
        data = response.json()
        if data['code'] != 200:
            return []
        data= data['data']

        new_data = []
        for item in data:
            i = item[0]    

            patent_id = i['id']
            sql = f"SELECT COUNT(*) FROM knowledgebase WHERE mark_info='{patent_id}' AND type_id = 2"
            id_result = kyzs_sql.mysql_exec(sql)[0]['COUNT(*)']

            i['is_collected'] = 1 if id_result else 0

            i['app_date'] = datetime.datetime.fromtimestamp(i['app_date']['seconds'])
            i['pub_date'] = datetime.datetime.fromtimestamp(i['pub_date']['seconds'])

            new_data.append(i)
        return new_data
    except requests.RequestException as e:
        print(f"请求出错: {e}")
        return None

def get_online_infomation_api(keyword:str):
    """
    获取讯飞互联网信息信息
    :param keyword: 关键词
    :return: 在线信息
    """
    def xunfei_online_search(question):
        url = "https://spark-api-open.xf-yun.com/v1/chat/completions"
        print(f'方法内：讯飞网络知识检索关键词：{question}', flush=True)
        tools = [
            {
                "type": "function",
                "function": {
                    "name": "get_research_information",
                    "description": "获取互联网科研信息",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "title": {
                                "type": "string",
                                "description": "信息标题"
                            },
                            "content": {
                                "type": "string",
                                "description": "信息内容，约1000字，倾向于真实案例和数据"
                            },
                            "source": {
                                "type": "string",
                                "description": "信息来源"
                            },
                            "date": {
                                "type": "string",
                                "description": "信息日期，格式为YYYY-MM-DD"
                            },
                            "author": {
                                "type": "string",
                                "description": "信息作者"
                            }
                        },
                        "required": ["title", "content", "source", "date", "author"]
                    }
                }
            },
            {
                "type": "web_search",
                "web_search": {
                    "enable": True,
                    "show_ref_label": True,
                    "search_mode": "deep"
                }
            }
        ]

        data = {
            "model": "4.0Ultra",
            "user": "default_user",
            "messages": [
                {
                "role": "system",
                "content": "你是知识渊博的助理，能够获取互联网科研信息"
            },
                {
                    "role": "user",
                    "content": question
                }
            ],
            "temperature": 1,
            "top_k": 6,
            "stream": False,
            "max_tokens": 5000,
            "tools": tools,
            "tool_choice": {
                "type": "function",
                "function": {
                    "name": "get_research_information"
                }
            },

        }
        header = {"Authorization": "Bearer fXMOutpFGqmwUoTTSpHz:OTsqdwPkqNmvlKJouKSU"}
        print(f'准备执行请求：{url}，请求头：{header}，请求数据：{data}', flush=True)
        try:
            response = requests.post(url, headers=header, json=data)
            print(f'请求已发送，等待响应...', flush=True)
        except Exception as e:
            print(f'发送请求时出错：{e}', flush=True)
            return None

        return response.json()


    return xunfei_online_search(keyword)


#重庆聚合文献接口
def get_article_from_juhe_api(keywords:list,date:str=None,page=1,size=10,sort=3):
    """
    从重庆聚合文献接口获取文献
    """

    def md5_encrypt(text: str) -> str:
        md5 = hashlib.md5()
        md5.update(text.encode('utf-8'))
        return md5.hexdigest()
    exps = []
    for keyword in keywords:
        exps.append(f"关键词:{keyword}")
    if date:
        exps.append(f"年份:{date}")
    exp = " AND ".join(exps)
    exp = urllib.parse.quote(exp)
    params = "{}|{}|{}|{}".format(exp, page, size, sort)
    print('表达式：',params)
    times = "{}".format(int(time.time() * 1000))
    en_params = base64.b64encode(params.encode('utf-8')).decode('utf-8')
    sign = "{}|{}|{}".format(en_params, "AF85101E523744ADA233DF14CCC76980", times)
    en_sign = md5_encrypt(sign)
    headers = {
        "User-Agent": "test",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection": "keep-alive",
    }
    query_params = {
        "params": en_params,
        "sign": en_sign,
        "token": "AF85101E523744ADA233DF14CCC76980",
        "times": times
    }
    try:
        with requests.Session() as session:
            url = "http://61.128.134.70:6655/groups/ask/search"
            session.get(url, headers=headers, params=query_params, timeout=10)
            response = session.get(url, headers=headers, params=query_params, timeout=10)

            print("Raw Response:")
            print(response.text)

            response.raise_for_status()
            result = response.json()['result']
            paperlist = []
            for record in result['records']:
                paperlist.append(
                    {
                        "标题": f"{record['title']}",
                        "关键词": f"{record['keyword']}",
                        "年份": f"{record['year']}",
                        "摘要": f"{record['content']}",
                        "DOI": f"{record['doi']}",
                        "下载链接": f"http://61.128.134.70:6655{record['full_source']}"
                    }
                )
            return paperlist
    except requests.exceptions.RequestException as e:
        print(f"请求发生错误: {str(e)}")
    except ValueError as e:
        print(f"JSON解析失败，请检查响应内容是否为有效JSON: {str(e)}")
    except KeyError as e:
        print(f"响应数据缺少必要字段: {str(e)}")
    except Exception as e:
        print(f"未知错误: {str(e)}")


#钻井院万方接口
import requests
import xml.etree.ElementTree as ET
def get_article_from_wanfang_api(Datewithin = None, Keywords = [], StartRecord = 1, MaximumRecords = 10):
    exps = []

    for keyword in Keywords:
        exps.append(f"{keyword}")

    if Datewithin:
        exps.append(f'Date within "{Datewithin}-01-01 {Datewithin}-12-31"')
    exp = " and ".join(exps)

    print(f"检索表达式：{exp}")
    def send_soap_request_paper(exp, startRecord = 1, maximumRecords = 10):
        url = "http://10.68.16.2/S/SRW/Paper.asmx"
        headers = {
            "Host": "10.68.16.2",
            "Content-Type": "text/xml; charset=utf-8",
            "SOAPAction": "searchRetrieve"
        }
        
        soap_body = f"""<?xml version="1.0" encoding="utf-8"?>
                    <soap:Envelope xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance" 
                                xmlns:xsd="http://www.w3.org/2001/XMLSchema" 
                                xmlns:soap="http://schemas.xmlsoap.org/soap/envelope/">
                    <soap:Body>
                        <searchRetrieveRequest xmlns="http://www.loc.gov/zing/srw/">
                            <version>1.2</version>
                            <query>{exp}</query>
                            <operation>searchretrieve</operation>
                            <startRecord>{startRecord}</startRecord>
                            <maximumRecords>{maximumRecords}</maximumRecords>
                        </searchRetrieveRequest>
                    </soap:Body>
                    </soap:Envelope>"""
        try:
            response = requests.post(url, data=soap_body, headers=headers)
            response.raise_for_status()
            print("Response Status Code:", response.status_code)
            return response.text
        except requests.exceptions.RequestException as e:
            print("Error sending SOAP request:", e)
            return None
    response = send_soap_request_paper(exp, startRecord=StartRecord, maximumRecords=MaximumRecords)
    with open("raw_response.txt", "w") as f:
        f.write(f"{response}")
    def parse_academic_papers(soap_response):
        namespaces = {
            "soap": "http://schemas.xmlsoap.org/soap/envelope/",
            "srw": "http://www.loc.gov/zing/srw/",
            "dc": "info:srw/schema/1/dc-v1.1",
            "srw_dc": "info:srw/schema/1/dc-v1.1"
        }
        print(soap_response)
        root = ET.fromstring(soap_response)
        number_element = root.find('.//srw:numberOfRecords', namespaces=namespaces)

        if number_element is not None:
            number_of_records = number_element.text
            print(f"共解析到论文数量: {number_of_records}")
        else:
            print("未找到numberOfRecords元素")
        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record",
            namespaces=namespaces
        )

        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record",
            namespaces=namespaces
        )

        papers = []
        for idx, record in enumerate(record_nodes, 1):
            dc_node = record.find("srw:recordData/srw_dc:dc", namespaces=namespaces)
            paper_info = {
                "标题": dc_node.find("dc:title", namespaces).text.strip() if dc_node.find("dc:title", namespaces).text is not None else "",
                "关键词": dc_node.find("dc:Subject", namespaces).text.strip() if dc_node.find("dc:Subject", namespaces).text is not None else "",
                "摘要": dc_node.find("dc:Description", namespaces).text.strip() if dc_node.find("dc:Description", namespaces).text is not None else "",
                "发表时间": dc_node.find("dc:Date", namespaces).text.strip() if dc_node.find("dc:Date", namespaces).text is not None else "",
                "DOI": dc_node.find("dc:Identifier", namespaces).text.strip() if dc_node.find("dc:Identifier", namespaces).text is not None else "",
            }
            paper_info["发表时间"] = paper_info["发表时间"].split("-")[0]
            papers.append(paper_info)

        return (papers, number_of_records)

    return parse_academic_papers(response)


#钻井院万方专利接口
def wangfang_patent(Datewithin = None, patent_name = [], StartRecord = 1, MaximumRecords = 10):
    exps = []

    for keyword in patent_name:
        exps.append(f"{keyword}")
        
    if Datewithin:
        exps.append(f'F_PublicationDate within "{Datewithin}-01-01 {Datewithin}-12-31"')
    exp = " and ".join(exps)

    def send_soap_request_patent(exp, startRecord = 1, maximumRecords = 10):
        url = "http://10.68.16.2/S/SRW/Patent.asmx"
        headers = {
            "Host": "10.68.16.2",
            "Content-Type": "text/xml; charset=utf-8",
            "SOAPAction": "searchRetrieve"
        }
        
        soap_body = f"""<?xml version="1.0" encoding="utf-8"?>
                    <soap:Envelope xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance" 
                                xmlns:xsd="http://www.w3.org/2001/XMLSchema" 
                                xmlns:soap="http://schemas.xmlsoap.org/soap/envelope/">
                    <soap:Body>
                        <searchRetrieveRequest xmlns="http://www.loc.gov/zing/srw/">
                            <version>1.2</version>
                            <query>{exp}</query>
                            <operation>searchretrieve</operation>
                            <startRecord>{startRecord}</startRecord>
                            <maximumRecords>{maximumRecords}</maximumRecords>
                        </searchRetrieveRequest>
                    </soap:Body>
                    </soap:Envelope>"""
        try:
            response = requests.post(url, data=soap_body, headers=headers)
            response.raise_for_status()
            print("Response Status Code:", response.status_code)
            return response.text
        except requests.exceptions.RequestException as e:
            print("Error sending SOAP request:", e)
            return None

    response = send_soap_request_patent(exp, startRecord=StartRecord, maximumRecords=MaximumRecords)

    def parse_patents(soap_response) -> tuple:
        namespaces = {
            "soap": "http://schemas.xmlsoap.org/soap/envelope/",
            "srw": "http://www.loc.gov/zing/srw/",
            "dc": "info:srw/schema/1/dc-v1.1",
            "srw_dc": "info:srw/schema/1/dc-v1.1"
        }
        root = ET.fromstring(soap_response)
        number_element = root.find('.//srw:numberOfRecords', namespaces=namespaces)

        if number_element is not None:
            number_of_records = number_element.text
            print(f"共解析到论文数量: {number_of_records}")
        else:
            print("未找到numberOfRecords元素")
        record_nodes = root.findall(
            ".//soap:Body/srw:searchRetrieveResponse/srw:records/srw:record",
            namespaces=namespaces
        )

        papers = []
        for idx, record in enumerate(record_nodes, 1):
            dc_node = record.find("srw:recordData/srw_dc:dc", namespaces=namespaces)

            paper_info = {
                "申请号": dc_node.find("dc:ApplicationNo", namespaces).text.strip() if dc_node.find("dc:ApplicationNo", namespaces).text is not None else "",
                "申请日": dc_node.find("dc:ApplicationDate", namespaces).text.strip() if dc_node.find("dc:ApplicationDate", namespaces).text is not None else "",
                "公开号": dc_node.find("dc:PublicationNo", namespaces).text.strip() if dc_node.find("dc:PublicationNo", namespaces).text is not None else "",
                "公开日": dc_node.find("dc:PublicationDate", namespaces).text.strip() if dc_node.find("dc:PublicationDate", namespaces).text is not None else "",
                "专利名称": dc_node.find("dc:PatentName", namespaces).text.strip() if dc_node.find("dc:PatentName", namespaces).text is not None else "",
                "申请人": dc_node.find("dc:Applicant", namespaces).text.strip() if dc_node.find("dc:Applicant", namespaces).text is not None else "",
                "发明人": dc_node.find("dc:Inventor", namespaces).text.strip() if dc_node.find("dc:Inventor", namespaces).text is not None else "",
                "IPC分类号": dc_node.find("dc:ClassMain", namespaces).text.strip() if dc_node.find("dc:ClassMain", namespaces).text is not None else "",
                "摘要": dc_node.find("dc:Abstract", namespaces).text.strip() if dc_node.find("dc:Abstract", namespaces).text is not None else "",
                "权利要求": dc_node.find("dc:SignoryItem", namespaces).text.strip() if dc_node.find("dc:SignoryItem", namespaces).text is not None else "",
            }
            papers.append(paper_info)

        return (papers, number_of_records)

    return parse_patents(response)
    

    
#大模型方法
def siliconflow_deepseek_answer(question):
    """
    deepseek大模型问答接口
    :param question: 题目
    :return:  回复文本（字符串）
    """
    url = "https://api.siliconflow.cn/v1/chat/completions"
    payload = {
        "model": "Pro/deepseek-ai/DeepSeek-V3",
        "messages": [
            {
                "role": "user",
                "content": str(question)
            }
        ],
        "stream": False,
        "max_tokens": 8192,
        "stop": ["null"],
        "temperature": 0.7,
        "top_p": 0.7,
        "top_k": 50,
        "frequency_penalty": 0.5,
        "n": 1,
        "response_format": {"type": "text"},
        "tools": [
            {
                "type": "function",
                "function": {
                    "description": "<string>",
                    "name": "<string>",
                    "parameters": {},
                    "strict": False
                }
            }
        ]
    }
    headers = {
        "Authorization": "Bearer sk-wmsgbfgsvjxjmyopswmaqfxnwtgmvtwqgsigehxmgwoihgeg",
        "Content-Type": "application/json"
    }

    response = requests.request("POST", url, json=payload, headers=headers)
    response = response.json()

    print(f"deepseek大模型返回结果：{response['choices'][0]['message']['content']}")
    return response['choices'][0]['message']['content']


def add_article_to_knowledge(article_data: dict,label_id:int,user_id:int):
    """
    文献收藏接口，将传入的文献数据插入到 SQL 数据库中
    """
    title = str(article_data.get('title_zh', article_data['title'])).replace("'", "").replace("\n","")
    abstract = article_data.get('abstract_zh', article_data.get('abstract', ''))
    keywords = article_data.get('keywords_zh',article_data['keywords'])

    content = '''
    文献信息如下:
    标题：{title}
    摘要：{abstract}
    关键词：{keywords}
    '''.format(title=title,abstract=abstract,keywords=keywords).replace("'", "").replace("\n","")

    mark_info = article_data.get('doi')
    sql_sentence =  f"""
    INSERT INTO `knowledgebase` (`title`, `content`, `label_id`, `user_id`, `type_id`, `mark_info`)
    VALUES ('{title}', '{content}', {label_id}, {user_id}, 1, '{mark_info}');
    """
    kyzs_sql.mysql_exec(sql_sentence)

    return {'code': 200, 'msg': 'success', 'data': None}


def add_patent_to_knowledge(patent_data: dict,label_id:int,user_id:int):
    """
    收藏专利
    """
    id = patent_data.get('id')
    title = str(patent_data.get('title')).replace("'", "\"")
    abstract = str(patent_data.get('abstract')).replace("'", "\"")
    claims = str(patent_data.get('claims')).replace("'", "\"")
    country = str(patent_data.get('country')).replace("'", "")
    app_num = str(patent_data.get('app_num')).replace("'", "")
    app_date = datetime.datetime.strptime(str(patent_data.get('app_date')), '%Y-%m-%dT%H:%M:%S')       
    pub_num = str(patent_data.get('pub_num')).replace("'", "")
    pub_date = datetime.datetime.strptime(str(patent_data.get('pub_date')), '%Y-%m-%dT%H:%M:%S')
    pub_kind = str(patent_data.get('pub_kind')).replace("'", "")
    applicant = str(patent_data.get('applicant')).replace("'", "\"")
    assignee = str(patent_data.get('assignee')).replace("'", "\"")
    inventor = str(patent_data.get('inventor')).replace("'", "\"")
    ipcr = str(patent_data.get('ipcr')).replace("'", "\"")       

    content = f'''
    专利摘要：{abstract}
    专利国家：{country}
    专利申请号：{app_num}
    专利申请日期：{app_date}
    专利公开号：{pub_num}
    专利公开日期：{pub_date}
    专利公开类型：{pub_kind}
    专利申请人：{applicant}
    '''
    sql=f"INSERT INTO knowledgebase (title, content, label_id, user_id, type_id,mark_info) VALUES ('{title}', '{content}', '{label_id}', '{user_id}', 2,'{id}')"
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def add_online_infomation_to_knowledge(online_infomation:dict,label_id:int,user_id:int):
    """
    收藏互联网信息
    """
    content= f'''
    日期：{str(online_infomation['date'])}
    标题：{str(online_infomation['title'])}
    内容：{str(online_infomation['content'])}
    来源：{str(online_infomation['source'])}
    '''
    sql=f"INSERT INTO knowledgebase (title, content, label_id, user_id, type_id,mark_info) VALUES ('{online_infomation['title']}', '{content}', '{label_id}', '{user_id}', 3,'网络信息收藏')"
    print("\033[31m" + "收藏互联网信息："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def add_mycontent_to_knowledge(content:str,title:str,label_id:int,user_id:int):
    """
    添加用户自定义内容
    """
    sql=f"INSERT INTO knowledgebase (title,content,label_id,user_id,type_id,mark_info) VALUES ('{title}','{content}',{label_id},{user_id},4,'用户自定义上传信息')"
    print("\033[31m" + "添加用户自定义内容："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def add_mycontent_file_to_knowledge(file_base64_string:str,file_extension:str,title:str,label_id:int,user_id:int):
    """
    添加用户自定义文件内容，读取文件base64，调用方法转为字符串，执行数据库插入
    """
    def text_extraction(b64_string, file_type):
        print(f'文件内容:{b64_string[:100]},文件类型:{file_type}')       
        import base64
        import os
        try:
            binary = base64.b64decode(b64_string)
        except:
            return "Invalid Base64 string", ""
        
        import time
        file_name = f"{int(time.time())}.{file_type}"
        
        def pdf_handler(file_binary):
            from pypdf import PdfReader
            with open("temp.pdf", "wb") as f:
                f.write(file_binary)
            reader = PdfReader("temp.pdf")
            full_text = ""
            for page in reader.pages:
                full_text += page.extract_text()
            os.remove("temp.pdf")
            return full_text, f"{int(time.time())}.pdf"
        
        def docx_handler(file_binary):
            from docx import Document
            with open("temp.docx", "wb") as f:
                f.write(file_binary)
            document = Document('temp.docx')
            full_text = ""
            for paragraph in document.paragraphs:
                full_text += paragraph.text
            os.remove("temp.docx")
            return full_text, f"{int(time.time())}.docx"
                    
        def txt_handler(file_binary):
            return file_binary.decode('utf-8'), f"{int(time.time())}.txt"
        
        def ppt_handler(file_binary):
            from pptx import Presentation
            with open("temp.pptx", "wb") as f:
                f.write(file_binary)
            prs = Presentation('temp.pptx')
            full_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
            os.remove("temp.pptx")
            return full_text, f"{int(time.time())}.pptx"
            
        content_string = ""
        match file_type:
            case "txt":
                content_string, file_name = txt_handler(binary)
            case "docx":
                content_string, file_name = docx_handler(binary)
            case "pdf":
                content_string, file_name = pdf_handler(binary)
            case "pptx":
                content_string, file_name = ppt_handler(binary)
            case _:
                file_name = f"{int(time.time())}.{file_type}"
 
        return content_string, file_name
    

    print(f'传入base64{file_base64_string[:100]}')
    sql=f"SELECT * FROM knowledgebase WHERE title = '{title}' and user_id = {user_id}"
    if kyzs_sql.mysql_exec(sql):
        return {"code": 400, "msg": "数据库当前用户中已存在相同文件名，不执行插入操作", "data": None}
    content_string, file_name = text_extraction(file_base64_string, str(file_extension).replace(".",""))
    content_string = content_string.replace("'", "''")
    content_string = content_string.replace("\n", "")
    real_title = title
    sql=f"INSERT INTO knowledgebase (title,content,label_id,user_id,type_id,mark_info) VALUES ('{real_title}','{content_string}',{label_id},{user_id},5,'用户自定义上传文件转换文本内容')"
    print("\033[31m" + "添加用户自定义文件内容："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

    
def modify_summary_api(review_id:int, start_position:int, end_position:int, replaced_text:str):
    sql_sentence = f"""
        SELECT review_body FROM review_records
        WHERE id = {review_id}
    """
    original_text = kyzs_sql.mysql_exec(sql_sentence)[0]['review_body']

    new_text = original_text[:start_position] + replaced_text + original_text[end_position:]

    sql_sentence = f"""
        UPDATE review_records
        SET review_body = '{new_text}'
        WHERE id = {review_id}
    """
    kyzs_sql.mysql_exec(sql_sentence)
    return 1

def modify_review_new_api(review_id:int, review_body:str):
    sql_sentence = f"""
        UPDATE review_records
        SET review_body = '{review_body}'
        WHERE id = {review_id}
    """
    kyzs_sql.mysql_exec(sql_sentence)
    return 1


def delete_summary(id:int):
    sql_sentence = f"""
        DELETE FROM review_records
        WHERE id = {id}
    """
    kyzs_sql.mysql_exec(sql_sentence)
    return 1


def add_format_api(format:dict):
    sql=f"INSERT INTO format (class_name, class_text) VALUES ('{format['class_name']}', '{format['class_text']}')"
    print("\033[31m" + "添加格式要求："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}

def delete_format_api(id:int):
    sql=f"DELETE FROM format WHERE id={id}"
    print("\033[31m" + "删除格式要求："+sql + "\033[0m")
    kyzs_sql.mysql_exec(sql)
    return {'code': 200, 'msg': 'success', 'data': None}


def generate_word_api(id:int):
    """
    传入综述id，生成Word文档
    """
    sql=f"SELECT * FROM review_records WHERE id={id}"
    print("\033[31m" + "获取综述内容用于生成word："+sql + "\033[0m")
    data = kyzs_sql.mysql_exec(sql)
    if data:
        md_text = data[0]['review_body']
    else:
        return {'code': 400, 'msg': '未找到该综述', 'data': None}

    md_text = md_text.replace('\\n', '\n')
    md_text = md_text.replace('```markdown', '')
    md_text = md_text.replace('```', '')
    
    md_file = 'temp.md'
    with open(md_file, 'w', encoding='utf-8') as f:
        f.write(md_text)
    
    output_path = "temp_output_word.docx"
    pypandoc.convert_file(md_file, to='docx', outputfile=output_path)
    with open(output_path, 'rb') as f:
        word_binary = f.read()

    os.remove(md_file)
    os.remove(output_path)

    import base64
    word_64 = base64.b64encode(word_binary).decode('utf-8')
    print('生成word文档成功')
    return {'code': 200, 'msg': 'success', 'data': word_64}


def generate_fuwenben_word_api(id: int):
    """
    传入综述id，生成富文本Word文档
    """
    sql = f"SELECT * FROM review_records WHERE id={id}"
    print("\033[31m" + "获取综述内容用于生成word：" + sql + "\033[0m")
    data = kyzs_sql.mysql_exec(sql)
    if not data:
        return {'code': 400, 'msg': '未找到该综述', 'data': None}
    
    html_text = data[0]['review_body']

    from docx import Document
    from bs4 import BeautifulSoup
    from io import BytesIO
    import base64

    doc = Document()
    
    soup = BeautifulSoup(html_text, 'html.parser')
    
    def process_element(element, parent_paragraph=None):
        if element.name is None:
            text = str(element)
            if parent_paragraph:
                run = parent_paragraph.add_run(text)
                return run
            return None
        
        if element.name in ['p', 'div', 'br']:
            if parent_paragraph and parent_paragraph.runs:
                pass
            new_p = doc.add_paragraph()
            for child in element.children:
                process_element(child, new_p)
            return None
        
        elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            level = int(element.name[1])
            heading = doc.add_heading(element.get_text(strip=False), level=level)
            return None
        
        elif element.name == 'strong' or element.name == 'b':
            run = None
            if parent_paragraph:
                run = parent_paragraph.add_run(element.get_text())
                run.bold = True
            for child in element.children:
                child_run = process_element(child, parent_paragraph)
                if child_run:
                    child_run.bold = True
            return run
        
        elif element.name == 'em' or element.name == 'i':
            run = None
            if parent_paragraph:
                run = parent_paragraph.add_run(element.get_text())
                run.italic = True
            for child in element.children:
                child_run = process_element(child, parent_paragraph)
                if child_run:
                    child_run.italic = True
            return run
        
        else:
            for child in element.children:
                process_element(child, parent_paragraph)
            return None

    for child in soup.children:
        process_element(child)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    
    word_b64 = base64.b64encode(buffer.read()).decode('utf-8')
    
    return {'code': 200, 'msg': 'success', 'data': word_b64}
